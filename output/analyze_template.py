import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import json

prs = Presentation(r'D:\BETEKK\opencode-config-template\opencode_app\scripts\templates\template.pptx')

for i, slide in enumerate(prs.slides):
    print(f'\n=== Slide {i} ===')
    layout = slide.slide_layout
    print(f'  Layout: {layout.name}')
    for shape in slide.shapes:
        stype = str(shape.shape_type)
        print(f'  Shape: {stype}, name={shape.name}')
        print(f'    pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})')
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                runs_text = ' '.join([r.text for r in para.runs])
                if runs_text.strip():
                    font_info = ''
                    if para.runs:
                        r = para.runs[0]
                        font_info = f' [font={r.font.name}, size={r.font.size}, bold={r.font.bold}]'
                    print(f'    Text: "{runs_text[:100]}"{font_info}')
        if shape.has_table:
            table = shape.table
            print(f'    Table: {len(list(table.rows))} rows x {len(table.columns)} cols')
            for row_idx, row in enumerate(table.rows):
                row_text = [cell.text for cell in row.cells]
                print(f'      Row {row_idx}: {row_text}')
