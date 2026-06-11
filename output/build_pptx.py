"""
build_pptx.py
=============
Produce a Chinese-localized version of the core BETEKK CanvasTEKK pitch deck
(slides 0-14) from template.pptx, dropping the appendix (slides 15-41).

Localizes the main slide titles and key callouts to Chinese while keeping
the brand name (BETEKK / CanvasTEKK). Targets shapes by their actual names in
the current template.

Usage:
    python output/build_pptx.py
"""

import sys
from pathlib import Path

from pptx import Presentation

sys.stdout.reconfigure(encoding="utf-8")

TEMPLATE = Path(__file__).resolve().parent.parent / "scripts" / "templates" / "template.pptx"
OUTPUT = Path(__file__).resolve().parent / "betekk_pitch_zh.pptx"

_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def delete_slide(prs, index: int) -> None:
    sld_id = prs.slides._sldIdLst[index]
    rid = sld_id.get(f"{{{_REL_NS}}}id")
    if rid:
        prs.part.drop_rel(rid)
    prs.slides._sldIdLst.remove(sld_id)


def find_shape(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_first_paragraph(shape, text: str) -> None:
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = text
        for run in tf.paragraphs[0].runs[1:]:
            run.text = ""
    elif tf.paragraphs:
        tf.paragraphs[0].text = text


# (slide_index, shape_name) -> Chinese text
L10N = {
    (0, "subtitle"): "建筑环境科技",
    (0, "tagline"): "每栋建筑都要接受成千上万次检查。我们让每一次检查数字化。",
    (1, "Title 3"): "建筑全生命周期中的每一次检查，至今仍依赖人工逐项完成",
    (2, "Title 4"): "设计审查、施工检查、交付核验——同一套人工流程贯穿三大阶段",
    (3, "title"): "这不仅是自动化工具，更是 92 亿美元的市场机遇",
    (3, "sec-head"): "市场布局",
    (4, "Title 13"): "三个阶段，数千次检查，一笔巨大的隐性成本",
    (5, "label-today-text"): "现状",
    (5, "Title 4"): "每一次建筑检查仍是人工完成，三个阶段概莫能外",
    (6, "Title 1"): "一个平台，覆盖每一次检查，由 CanvasTEKK 全自动完成",
    (7, "Title 1"): "构建任意检查工作流，一键运行，即时获取结论",
    (8, "Title 1"): "你的规则，任意模型，自动发现每一处差异",
    (9, "title"): "按使用量计费——每一次工作流运行，ROI 持续放大",
    (10, "title"): "CanvasTEKK 为何能赢，又为何难以复制",
    (11, "title"): "市场验证与进展",
    (12, "Text 0"): "核心团队",
    (13, "Title 14"): "变革建筑业，正当其时：对的团队、对的技术、对的时机",
    (14, "Title 1"): "感谢聆听",
    (14, "Subtitle 2"): "访问 https://betekk.com  |  CanvasTEKK 由 BETEKK 出品",
}


def main() -> None:
    if not TEMPLATE.exists():
        print(f"ERROR: template not found: {TEMPLATE}")
        sys.exit(1)

    prs = Presentation(str(TEMPLATE))
    total = len(prs.slides)
    print(f"Loaded template: {total} slides")

    # Keep slides 0-14 (main pitch); delete appendix 15..(total-1) from the top.
    for index in range(total - 1, 14, -1):
        delete_slide(prs, index)
    print(f"Trimmed to {len(prs.slides)} slides (main pitch)")

    slides = list(prs.slides)

    applied = 0
    for (slide_idx, shape_name), text in L10N.items():
        if slide_idx >= len(slides):
            continue
        shape = find_shape(slides[slide_idx], shape_name)
        if shape is not None:
            set_first_paragraph(shape, text)
            applied += 1
    print(f"Localized {applied}/{len(L10N)} text fields")

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")

    # Verify
    prs2 = Presentation(str(OUTPUT))
    print(f"\nVerification - {len(prs2.slides)} slides:")
    for i, slide in enumerate(prs2.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = "".join(r.text for r in shape.text_frame.paragraphs[0].runs).strip()
            if txt:
                print(f"  Slide {i}: {txt[:60]}")
                break


if __name__ == "__main__":
    main()
