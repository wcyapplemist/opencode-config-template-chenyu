import sys
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

TEMPLATE = r'D:\BETEKK\opencode-config-template\opencode_app\scripts\templates\template.pptx'
OUTPUT = r'D:\BETEKK\opencode-config-template\opencode_app\output\digital_construction.pptx'

prs = Presentation(TEMPLATE)

def delete_slide(prs, index):
    rId = prs.slides._sldIdLst[index].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    sldId = prs.slides._sldIdLst[index]
    prs.slides._sldIdLst.remove(sldId)

def replace_in_shape(shape, mapping):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for old, new in mapping.items():
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)

def set_para_texts(shape, texts):
    if not shape.has_text_frame:
        return
    for i, para in enumerate(shape.text_frame.paragraphs):
        if i < len(texts):
            if para.runs:
                para.runs[0].text = texts[i]
                for r in para.runs[1:]:
                    r.text = ""
            else:
                para.text = texts[i]
        else:
            for r in para.runs:
                r.text = ""

def set_table_data(table, data):
    for r_idx, row_data in enumerate(data):
        for c_idx, cell_text in enumerate(row_data):
            if r_idx < len(table.rows) and c_idx < len(table.columns):
                cell = table.cell(r_idx, c_idx)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
                if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                    cell.text_frame.paragraphs[0].runs[0].text = cell_text
                elif cell.text_frame.paragraphs:
                    cell.text_frame.paragraphs[0].text = cell_text

for idx in [10, 8, 5]:
    delete_slide(prs, idx)

# Remaining order (no swap needed):
# 0: Title Slide (Cover)
# 1: Agenda
# 2: Section Header (full image + centered title)
# 3: Section Header with subtitle (image + title + subtitle)
# 4: Title + content + image
# 5: Two Content dark (was NAVIGATING Q&A SESSIONS)
# 6: Two Content white (was EFFECTIVE DELIVERY TECHNIQUES)
# 7: Title + content + table (was DYNAMIC DELIVERY)
# 8: Two Content dark (was Final tips & takeaways)
# 9: Closing (was Thank you)

slides = list(prs.slides)
print(f"Total slides: {len(slides)}")

# ===== Slide 0: Cover =====
for shape in slides[0].shapes:
    if shape.name == "Title 3":
        set_para_texts(shape, ["数字化技术赋能建筑领域"])
    elif shape.name == "Subtitle 2":
        set_para_texts(shape, ["DIGITAL TECHNOLOGY IN CONSTRUCTION"])

# ===== Slide 1: Agenda =====
for shape in slides[1].shapes:
    if shape.name == "Title 1":
        replace_in_shape(shape, {"Agenda": "目录"})
    elif shape.name == "Subtitle 2":
        set_para_texts(shape, [
            "BIM 建筑信息模型",
            "智慧工地与物联网",
            "大数据与人工智能",
            "数字化转型效益分析",
            "未来展望与趋势"
        ])

# ===== Slide 2: Section Header → 数字化技术概览 =====
for shape in slides[2].shapes:
    if shape.name == "Title 5":
        replace_in_shape(shape, {"THE POWER OF COMMUNICATION": "数字化技术赋能建筑"})

# ===== Slide 3: Section Header with subtitle → BIM =====
for shape in slides[3].shapes:
    if shape.name == "Title 6":
        replace_in_shape(shape, {"OVERCOMING NERVOUSNESS": "BIM 建筑信息模型"})
    elif shape.name == "Text Placeholder 7":
        replace_in_shape(shape, {"Confidence-building strategies": "Building Information Modeling"})

# ===== Slide 4: Title + content + image → BIM核心应用 =====
for shape in slides[4].shapes:
    if shape.name == "Title 3":
        set_para_texts(shape, ["BIM 核心应用场景"])
    elif shape.name == "Subtitle 4":
        set_para_texts(shape, [
            "BIM 在建筑全生命周期的深度应用",
            "三维可视化协同设计，提升设计质量与协调效率",
            "碰撞检测与冲突分析，减少施工返工率达30%",
            "工程量自动统计，精确成本管控与预算优化",
            "4D施工进度模拟，优化工期安排与资源配置"
        ])

# ===== Slide 5: Two Content dark → 智慧工地与物联网 =====
for shape in slides[5].shapes:
    if shape.name == "Title 3":
        replace_in_shape(shape, {"NAVIGATING Q&A SESSIONS": "智慧工地与物联网"})
    elif shape.name == "Content Placeholder 7":
        set_para_texts(shape, [
            "智慧工地感知网络",
            "传感器实时采集施工环境与设备运行数据",
            "人员定位与电子围栏安全预警系统",
            "大型设备远程监控与预测性维护"
        ])
    elif shape.name == "Content Placeholder 8":
        set_para_texts(shape, [
            "物联网平台架构",
            "边缘计算网关实现数据本地实时处理",
            "云平台统一管理与跨项目数据协同",
            "5G网络保障工地高速低延迟通信"
        ])

# ===== Slide 6: Two Content white → AI与自动化 =====
for shape in slides[6].shapes:
    if shape.name == "Title 1":
        replace_in_shape(shape, {"EFFECTIVE DELIVERY TECHNIQUES": "大数据与人工智能"})
    elif shape.name == "Content Placeholder 2":
        set_para_texts(shape, [
            "大数据智能分析",
            "海量施工数据驱动科学决策与风险预判",
            "资源配置优化与供应链智能管理",
            "质量缺陷AI识别与自动追溯分析"
        ])
    elif shape.name == "Content Placeholder 3":
        set_para_texts(shape, [
            "AI 自动化技术",
            "建筑机器人执行高危与重复性施工作业",
            "3D打印技术实现复杂结构快速精准建造",
            "无人机自动化巡检与土方测量",
            "智能物流与自动化物料配送系统"
        ])

# ===== Slide 7: Title + content + table → 数字化效益 =====
for shape in slides[7].shapes:
    if shape.name == "Title 1":
        replace_in_shape(shape, {"DYNAMIC DELIVERY": "数字化转型效益分析"})
    elif shape.name == "Content Placeholder 2":
        set_para_texts(shape, [
            "数字化技术正在重塑建筑行业的生产与管理方式",
            "显著提升项目效率、工程质量与安全管理水平"
        ])
    elif shape.has_table:
        set_table_data(shape.table, [
            ["效益指标", "计量单位", "传统方式", "数字化方式"],
            ["设计协调效率", "提升百分比", "基准", "+40%"],
            ["施工返工率", "百分比 (%)", "12%", "3%"],
            ["安全事故率", "降低百分比", "基准", "-60%"],
            ["项目成本节约", "百分比 (%)", "基准", "+25%"],
            ["工期缩短比例", "百分比 (%)", "基准", "+20%"]
        ])

# ===== Slide 8: Two Content dark → 未来展望 =====
for shape in slides[8].shapes:
    if shape.name == "Title 1":
        replace_in_shape(shape, {"Final tips & takeaways": "未来展望与发展趋势"})
    elif shape.name == "Content Placeholder 9":
        set_para_texts(shape, [
            "技术发展前沿",
            "生成式AI辅助建筑创意设计",
            "数字孪生城市与建筑运维管理",
            "碳中和智慧建筑系统",
            "区块链建筑供应链透明管理"
        ])
    elif shape.name == "Content Placeholder 15":
        set_para_texts(shape, [
            "行业变革方向",
            "建筑工业化与装配式智能建造",
            "全流程数字化协同管理平台",
            "跨界融合创新生态构建",
            "绿色低碳可持续建筑发展",
            "人机协作智能建造新模式"
        ])

# ===== Slide 9: Closing =====
for shape in slides[9].shapes:
    if shape.name == "Title 2":
        replace_in_shape(shape, {"Thank you": "感谢聆听"})
    elif shape.name == "Subtitle 3":
        set_para_texts(shape, [
            "数字化赋能 · 建筑未来",
            "Digital Empowerment · Future Construction"
        ])

prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")

# Verify
prs2 = Presentation(OUTPUT)
print(f"\nVerification - {len(prs2.slides)} slides:")
for i, slide in enumerate(prs2.slides):
    t = ""
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs).strip()[:60]
                if t: break
        if t: break
    print(f"  Slide {i}: {t}")
