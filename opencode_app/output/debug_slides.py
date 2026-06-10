"""
debug_slides.py
===============
Generic slide inspector for template.pptx.

Dumps every slide's text shapes (name + text) so you can discover the exact
shape names needed to edit the BETEKK CanvasTEKK deck. Read-only — does not
modify or save the template.

Usage:
    python output/debug_slides.py
"""

import sys
from pathlib import Path

from pptx import Presentation

sys.stdout.reconfigure(encoding="utf-8")

TEMPLATE = Path(__file__).resolve().parent.parent / "scripts" / "templates" / "template.pptx"


def shape_text(shape) -> str:
    if not shape.has_text_frame:
        return "(no text frame)"
    parts = []
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text for run in para.runs).strip()
        if line:
            parts.append(line)
    return " | ".join(parts)


def main() -> None:
    if not TEMPLATE.exists():
        print(f"ERROR: template not found: {TEMPLATE}")
        sys.exit(1)

    prs = Presentation(str(TEMPLATE))
    print(f"Template: {TEMPLATE.name}")
    print(f"Slides: {len(prs.slides)}  |  Layouts: {len(prs.slide_layouts)}\n")

    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape_text(shape)
            if text:
                print(f"  [{shape.name}] {text[:80]}")
        print()


if __name__ == "__main__":
    main()
