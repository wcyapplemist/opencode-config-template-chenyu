import sys
from pptx import Presentation
sys.stdout.reconfigure(encoding='utf-8')

prs = Presentation(r'D:\BETEKK\opencode-config-template\opencode_app\output\digital_construction.pptx')
for i, slide in enumerate(prs.slides):
    print(f"\n{'='*50}")
    print(f"Slide {i+1}")
    print(f"{'='*50}")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = "".join(r.text for r in para.runs).strip()
                if txt:
                    print(f"  {txt}")
        if shape.has_table:
            print("  [TABLE]")
            for row in shape.table.rows:
                cells = [cell.text for cell in row.cells]
                print(f"  | {' | '.join(cells)} |")
