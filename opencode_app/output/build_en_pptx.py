import sys
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

TEMPLATE = r'D:\BETEKK\opencode-config-template\opencode_app\scripts\templates\template.pptx'
OUTPUT = r'D:\BETEKK\opencode-config-template\opencode_app\output\digital_construction_en.pptx'

prs = Presentation(TEMPLATE)

def delete_slide(prs, index):
    rId = prs.slides._sldIdLst[index].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    sldId = prs.slides._sldIdLst[index]
    prs.slides._sldIdLst.remove(sldId)

def set_para_texts(shape, texts):
    if not shape.has_text_frame:
        return
    for i, para in enumerate(shape.text_frame.paragraphs):
        if i < len(texts):
            if para.runs:
                para.runs[0].text = texts[i]
                for r in para.runs[1:]:
                    r.text = ""
            else:
                para.text = texts[i]
        else:
            for r in para.runs:
                r.text = ""

def replace_in_shape(shape, mapping):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for old, new in mapping.items():
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)

def set_table_data(table, data):
    for r_idx, row_data in enumerate(data):
        for c_idx, cell_text in enumerate(row_data):
            if r_idx < len(table.rows) and c_idx < len(table.columns):
                cell = table.cell(r_idx, c_idx)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
                if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                    cell.text_frame.paragraphs[0].runs[0].text = cell_text
                elif cell.text_frame.paragraphs:
                    cell.text_frame.paragraphs[0].text = cell_text

for idx in [10, 8, 5]:
    delete_slide(prs, idx)

slides = list(prs.slides)
print(f"Total slides: {len(slides)}")

# Slide 0: Cover
for shape in slides[0].shapes:
    if shape.name == "Title 3":
        set_para_texts(shape, ["Digital Technology Empowering Construction"])
    elif shape.name == "Subtitle 2":
        set_para_texts(shape, ["Transforming the Built Environment with Innovation"])

# Slide 1: Agenda
for shape in slides[1].shapes:
    if shape.name == "Title 1":
        replace_in_shape(shape, {"Agenda": "Agenda"})
    elif shape.name == "Subtitle 2":
        set_para_texts(shape, [
            "Introduction to Digital Construction",
            "BIM Technology",
            "Smart Construction Sites & IoT",
            "Big Data & AI in Construction",
            "Digital Transformation Benefits"
        ])

# Slide 2: Section Header - Introduction
for shape in slides[2].shapes:
    if shape.name == "Title 5":
        replace_in_shape(shape, {
            "THE POWER OF COMMUNICATION": "INTRODUCTION TO DIGITAL CONSTRUCTION"
        })

# Slide 3: Section Header with subtitle - BIM
for shape in slides[3].shapes:
    if shape.name == "Title 6":
        replace_in_shape(shape, {
            "OVERCOMING NERVOUSNESS": "BIM"
        })
    elif shape.name == "Text Placeholder 7":
        replace_in_shape(shape, {
            "Confidence-building strategies": "Building Information Modeling"
        })

# Slide 4: Title + content + image - BIM Applications
for shape in slides[4].shapes:
    if shape.name == "Title 3":
        set_para_texts(shape, ["BIM Core Applications"])
    elif shape.name == "Subtitle 4":
        set_para_texts(shape, [
            "Key applications across the building lifecycle",
            "3D Visualization for collaborative design and coordination",
            "Clash Detection to reduce rework by up to 30%",
            "Automated Quantity Takeoff for precise cost control",
            "4D Schedule Simulation to optimize project timelines"
        ])

# Slide 5: Two Content dark - Smart Construction & IoT
for shape in slides[5].shapes:
    if shape.name == "Title 3":
        replace_in_shape(shape, {
            "NAVIGATING Q&A SESSIONS": "SMART CONSTRUCTION & IoT"
        })
    elif shape.name == "Content Placeholder 7":
        set_para_texts(shape, [
            "IoT Sensor Networks",
            "Real-time environmental and equipment monitoring",
            "Worker location tracking and safety alerts",
            "Predictive maintenance for heavy machinery"
        ])
    elif shape.name == "Content Placeholder 8":
        set_para_texts(shape, [
            "Platform Architecture",
            "Edge computing for on-site data processing",
            "Cloud-based cross-project data collaboration",
            "5G-enabled high-speed low-latency communication"
        ])

# Slide 6: Two Content white - Big Data & AI
for shape in slides[6].shapes:
    if shape.name == "Title 1":
        replace_in_shape(shape, {
            "EFFECTIVE DELIVERY TECHNIQUES": "BIG DATA & AI IN CONSTRUCTION"
        })
    elif shape.name == "Content Placeholder 2":
        set_para_texts(shape, [
            "Big Data Analytics",
            "Data-driven decision making and risk prediction",
            "Resource optimization and supply chain management",
            "AI-powered defect detection and root cause analysis"
        ])
    elif shape.name == "Content Placeholder 3":
        set_para_texts(shape, [
            "AI Automation Technologies",
            "Construction robots for hazardous repetitive tasks",
            "3D printing for complex structures",
            "Drone-based automated inspection and surveying",
            "Smart logistics and automated material delivery"
        ])

# Slide 7: Title + content + table - Benefits
for shape in slides[7].shapes:
    if shape.name == "Title 1":
        replace_in_shape(shape, {"DYNAMIC DELIVERY": "DIGITAL TRANSFORMATION BENEFITS"})
    elif shape.name == "Content Placeholder 2":
        set_para_texts(shape, [
            "Digital technology is reshaping the construction industry",
            "Significant gains in efficiency, quality, and safety"
        ])
    elif shape.has_table:
        set_table_data(shape.table, [
            ["Metric", "Unit", "Traditional", "Digital"],
            ["Design Coordination", "Improvement", "Baseline", "+40%"],
            ["Construction Rework", "Percentage", "12%", "3%"],
            ["Safety Incidents", "Reduction", "Baseline", "-60%"],
            ["Cost Savings", "Percentage", "Baseline", "+25%"],
            ["Schedule Compression", "Percentage", "Baseline", "+20%"]
        ])

# Slide 8: Two Content dark - Future Trends
for shape in slides[8].shapes:
    if shape.name == "Title 1":
        replace_in_shape(shape, {
            "Final tips & takeaways": "FUTURE TRENDS & OUTLOOK"
        })
    elif shape.name == "Content Placeholder 9":
        set_para_texts(shape, [
            "Technology Frontiers",
            "Generative AI for architectural design",
            "Digital twin cities and building operations",
            "Carbon-neutral smart building systems",
            "Blockchain for supply chain transparency"
        ])
    elif shape.name == "Content Placeholder 15":
        set_para_texts(shape, [
            "Industry Transformation",
            "Industrialized and prefabricated construction",
            "End-to-end digital collaboration platforms",
            "Cross-sector innovation ecosystems",
            "Green and sustainable building development",
            "Human-robot collaborative construction"
        ])

# Slide 9: Closing
for shape in slides[9].shapes:
    if shape.name == "Title 2":
        replace_in_shape(shape, {"Thank you": "Thank You"})
    elif shape.name == "Subtitle 3":
        set_para_texts(shape, [
            "Digital Empowerment · Future Construction",
            "info@digitalconstruction.com"
        ])

prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")

# Verify
prs2 = Presentation(OUTPUT)
print(f"\nVerification - {len(prs2.slides)} slides:")
for i, slide in enumerate(prs2.slides):
    t = ""
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs).strip()[:70]
                if t: break
        if t: break
    print(f"  Slide {i}: {t}")
