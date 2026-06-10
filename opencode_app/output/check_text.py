import sys
from pptx import Presentation
sys.stdout.reconfigure(encoding='utf-8')

prs = Presentation(r'D:\BETEKK\opencode-config-template\opencode_app\scripts\templates\template.pptx')
slide = list(prs.slides)[4]
for shape in slide.shapes:
    if shape.name == "Title 3":
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                print(f"Run text repr: {repr(run.text)}")
