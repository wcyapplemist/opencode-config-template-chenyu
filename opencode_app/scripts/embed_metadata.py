"""
embed_metadata.py
=================
One-time script that injects <p:ext> layout metadata into a template .pptx.

Each Slide Master layout receives a custom XML block describing:
  - templateId: unique semantic identifier (e.g. "content_three_column")
  - label: human-readable name
  - useWhen: guidance for when to choose this layout
  - compatibleWith: maps back to existing ppt_builder slide_type (backward compat)
  - slots: per-placeholder role, required flag, and accepted content types

Usage:
    python scripts/embed_metadata.py [input.pptx] [output.pptx]

    Default input:  scripts/templates/template.pptx
    Default output: scripts/templates/template_tagged.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path
from typing import Any, Dict, List, Optional

from lxml import etree
from pptx import Presentation

_SCRIPT_DIR = Path(__file__).resolve().parent
TEMPLATES_DIR = _SCRIPT_DIR / "templates"

CUSTOM_NS = "https://beteek.com/pptx-layout-meta"
CUSTOM_URI = "{b9e7a3d1-4c5f-4a8b-b2d6-e1f3a7c9d0e2}"

_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

LAYOUT_METADATA: Dict[str, Dict[str, Any]] = {
    "Title Slide": {
        "templateId": "title_slide",
        "label": "Cover Title Slide",
        "useWhen": "First slide of any deck, presentation opener",
        "compatibleWith": "title_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 1, "name": "Subtitle 2", "role": "subtitle", "required": False, "accepts": "text"},
        ],
    },
    "Section Header": {
        "templateId": "section_header",
        "label": "Section Header",
        "useWhen": "Major section or chapter break with a subtitle area",
        "compatibleWith": "section_header_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 1, "name": "Text Placeholder 2", "role": "subtitle", "required": False, "accepts": "text"},
        ],
    },
    "1_Section Header": {
        "templateId": "section_header_with_content",
        "label": "Section Header with Side Content",
        "useWhen": "Section break that also shows two small content blocks on the right",
        "compatibleWith": "section_header_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 1, "name": "Text Placeholder 2", "role": "subtitle", "required": False, "accepts": "text"},
            {"idx": 13, "name": "Content Placeholder 9", "role": "side_content_top", "required": False, "accepts": "bullets|text"},
            {"idx": 14, "name": "Content Placeholder 11", "role": "side_content_bottom", "required": False, "accepts": "bullets|text"},
        ],
    },
    "1_Section": {
        "templateId": "section_plain",
        "label": "Plain Section Divider",
        "useWhen": "Minimal section break, full-width title with body text",
        "compatibleWith": "section_header_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 1, "name": "Text Placeholder 2", "role": "subtitle", "required": False, "accepts": "text"},
        ],
    },
    "Main title page": {
        "templateId": "title_with_summary",
        "label": "Title Page with Summary",
        "useWhen": "Cover slide that includes a short summary or tagline below the title",
        "compatibleWith": "title_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 11, "name": "Text Placeholder 4", "role": "subtitle", "required": False, "accepts": "text"},
            {"idx": 12, "name": "Text Placeholder 6", "role": "summary", "required": False, "accepts": "text"},
        ],
    },
    "Section": {
        "templateId": "section_minimal",
        "label": "Minimal Section Divider",
        "useWhen": "Clean section break with title and optional body text",
        "compatibleWith": "section_header_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 1, "name": "Text Placeholder 2", "role": "subtitle", "required": False, "accepts": "text"},
        ],
    },
    "1_Blank": {
        "templateId": "blank",
        "label": "Blank Layout",
        "useWhen": "Empty canvas for free-form content; not recommended for auto-generation",
        "compatibleWith": None,
        "slots": [],
    },
    "Title Only": {
        "templateId": "title_only",
        "label": "Title Only",
        "useWhen": "Slide with only a title bar, rest is blank for images or custom content",
        "compatibleWith": "section_header_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
        ],
    },
    "1_Title and Content": {
        "templateId": "content_full",
        "label": "Full-Width Content",
        "useWhen": "Standard content slide with title and one large body area",
        "compatibleWith": "content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 1, "name": "Content Placeholder 2", "role": "body", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "7_Title and Content": {
        "templateId": "content_with_subtitle",
        "label": "Content with Subtitle Bar",
        "useWhen": "Content slide with a subtitle/header bar above the main body",
        "compatibleWith": "content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 14, "name": "Text Placeholder 3", "role": "subtitle", "required": False, "accepts": "text"},
            {"idx": 1, "name": "Content Placeholder 2", "role": "body", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "Blank": {
        "templateId": "blank_with_title",
        "label": "Blank with Optional Title",
        "useWhen": "Near-blank slide with an optional title; rest is free-form",
        "compatibleWith": None,
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": False, "accepts": "text"},
        ],
    },
    "Title and Content": {
        "templateId": "content_two_column_mixed",
        "label": "Two Column Mixed (OBJECT + BODY)",
        "useWhen": "Two-column layout with one rich-content column and one text column",
        "compatibleWith": "two_content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 1, "name": "Content Placeholder 2", "role": "left_content", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 12, "name": "Text Placeholder 15", "role": "right_content", "required": True, "accepts": "bullets|text"},
        ],
    },
    "3_Title and Content": {
        "templateId": "content_two_object",
        "label": "Two Column Equal (OBJECT + OBJECT)",
        "useWhen": "Side-by-side comparison with two equal rich-content areas",
        "compatibleWith": "two_content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 1, "name": "Content Placeholder 2", "role": "left_content", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 15, "name": "Content Placeholder 2", "role": "right_content", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "2_Title and Content": {
        "templateId": "content_two_object_notitle",
        "label": "Two Column Without Title Row",
        "useWhen": "Two side-by-side content areas with a title but no separate title row",
        "compatibleWith": "two_content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 1, "name": "Content Placeholder 2", "role": "left_content", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 12, "name": "Content Placeholder 2", "role": "right_content", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "23_Title and Content": {
        "templateId": "content_two_object_with_header",
        "label": "Two Column with Header Bar",
        "useWhen": "Two rich-content columns with a subtitle/header bar spanning the top",
        "compatibleWith": "two_content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 14, "name": "Text Placeholder 3", "role": "subtitle", "required": False, "accepts": "text"},
            {"idx": 18, "name": "Content Placeholder 8", "role": "left_content", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 19, "name": "Content Placeholder 10", "role": "right_content", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "4_Title and Content": {
        "templateId": "content_two_object_titled",
        "label": "Two Column with Individual Subtitles",
        "useWhen": "Two content columns each with their own subtitle label above",
        "compatibleWith": "two_content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 14, "name": "Text Placeholder 3", "role": "left_subtitle", "required": False, "accepts": "text"},
            {"idx": 15, "name": "Text Placeholder 3", "role": "right_subtitle", "required": False, "accepts": "text"},
            {"idx": 18, "name": "Content Placeholder 8", "role": "left_content", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 19, "name": "Content Placeholder 10", "role": "right_content", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "21_Title and Content": {
        "templateId": "content_two_object_asymmetric",
        "label": "Two Column Asymmetric with Subtitle",
        "useWhen": "Two content columns where one is taller, with a subtitle over the left column",
        "compatibleWith": "two_content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 14, "name": "Text Placeholder 3", "role": "left_subtitle", "required": False, "accepts": "text"},
            {"idx": 18, "name": "Content Placeholder 8", "role": "left_content", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 19, "name": "Content Placeholder 10", "role": "right_content", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "22_Title and Content": {
        "templateId": "content_two_object_alt",
        "label": "Two Column Alternate Layout",
        "useWhen": "Two content columns, right column has its own subtitle",
        "compatibleWith": "two_content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 15, "name": "Text Placeholder 3", "role": "right_subtitle", "required": False, "accepts": "text"},
            {"idx": 18, "name": "Content Placeholder 8", "role": "left_content", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 19, "name": "Content Placeholder 10", "role": "right_content", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "11_Title and Content": {
        "templateId": "content_stacked_two",
        "label": "Stacked Two Content Areas",
        "useWhen": "Two content areas stacked vertically (top and bottom)",
        "compatibleWith": "content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 1, "name": "Content Placeholder 2", "role": "body_top", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 14, "name": "Content Placeholder 2", "role": "body_bottom", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "24_Title and Content": {
        "templateId": "content_stacked_with_header",
        "label": "Stacked with Header Bar",
        "useWhen": "Three-tier layout: subtitle bar, top content, bottom content",
        "compatibleWith": "content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 15, "name": "Text Placeholder 3", "role": "subtitle", "required": False, "accepts": "text"},
            {"idx": 1, "name": "Content Placeholder 2", "role": "body_top", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 14, "name": "Content Placeholder 2", "role": "body_bottom", "required": False, "accepts": "bullets|text|richtext"},
        ],
    },
    "10_Title and Content": {
        "templateId": "content_stacked_equal",
        "label": "Stacked Equal Halves",
        "useWhen": "Two equally-sized content areas stacked vertically",
        "compatibleWith": "content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 1, "name": "Content Placeholder 2", "role": "body_top", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 12, "name": "Content Placeholder 2", "role": "body_bottom", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "5_Title and Content": {
        "templateId": "content_three_column",
        "label": "Three Column Content",
        "useWhen": "Three parallel topics, three features, three-phase comparison",
        "compatibleWith": None,
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 20, "name": "Content Placeholder 12", "role": "col_1", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 18, "name": "Content Placeholder 6", "role": "col_2", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 19, "name": "Content Placeholder 10", "role": "col_3", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "14_Title and Content": {
        "templateId": "content_three_with_header",
        "label": "Three Column with Header Bar",
        "useWhen": "Three content columns with a subtitle/header spanning the top",
        "compatibleWith": None,
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 14, "name": "Text Placeholder 3", "role": "subtitle", "required": False, "accepts": "text"},
            {"idx": 19, "name": "Content Placeholder 11", "role": "col_1", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 20, "name": "Content Placeholder 13", "role": "col_2", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 21, "name": "Content Placeholder 15", "role": "col_3", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "20_Title and Content": {
        "templateId": "content_grid_2x2",
        "label": "2x2 Grid Layout",
        "useWhen": "Four content blocks in a grid, four-quadrant analysis",
        "compatibleWith": None,
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 1, "name": "Content Placeholder 2", "role": "cell_top_left", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 17, "name": "Content Placeholder 2", "role": "cell_top_right", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 16, "name": "Content Placeholder 2", "role": "cell_bottom_left", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 18, "name": "Content Placeholder 2", "role": "cell_bottom_right", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "9_Title and Content": {
        "templateId": "content_stacked_with_labels",
        "label": "Stacked with Individual Labels",
        "useWhen": "Two stacked content areas, each with its own label/header",
        "compatibleWith": "content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 14, "name": "Text Placeholder 3", "role": "top_label", "required": False, "accepts": "text"},
            {"idx": 15, "name": "Text Placeholder 3", "role": "bottom_label", "required": False, "accepts": "text"},
            {"idx": 1, "name": "Content Placeholder 2", "role": "body_top", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 12, "name": "Content Placeholder 2", "role": "body_bottom", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "8_Title and Content": {
        "templateId": "gallery_6_images",
        "label": "Six Image Gallery",
        "useWhen": "Image gallery with six picture placeholders in asymmetric grid",
        "compatibleWith": None,
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 16, "name": "Picture Placeholder 5", "role": "image_1", "required": False, "accepts": "image"},
            {"idx": 12, "name": "Picture Placeholder 5", "role": "image_2", "required": False, "accepts": "image"},
            {"idx": 15, "name": "Picture Placeholder 5", "role": "image_3", "required": False, "accepts": "image"},
            {"idx": 17, "name": "Picture Placeholder 5", "role": "image_4", "required": False, "accepts": "image"},
            {"idx": 18, "name": "Picture Placeholder 5", "role": "image_5", "required": False, "accepts": "image"},
            {"idx": 19, "name": "Picture Placeholder 5", "role": "image_6", "required": False, "accepts": "image"},
        ],
    },
    "15_Title and Content": {
        "templateId": "gallery_3_text3_images",
        "label": "Three Column Text-Image Mixed",
        "useWhen": "Three columns each with a text label on top and image below",
        "compatibleWith": None,
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 14, "name": "Text Placeholder 3", "role": "col_1_label", "required": False, "accepts": "text"},
            {"idx": 19, "name": "Text Placeholder 3", "role": "col_2_label", "required": False, "accepts": "text"},
            {"idx": 20, "name": "Text Placeholder 3", "role": "col_3_label", "required": False, "accepts": "text"},
            {"idx": 21, "name": "Picture Placeholder 5", "role": "col_1_image", "required": False, "accepts": "image"},
            {"idx": 22, "name": "Picture Placeholder 5", "role": "col_2_image", "required": False, "accepts": "image"},
            {"idx": 23, "name": "Picture Placeholder 5", "role": "col_3_image", "required": False, "accepts": "image"},
            {"idx": 16, "name": "Picture Placeholder 5", "role": "col_1_image_bottom", "required": False, "accepts": "image"},
            {"idx": 17, "name": "Picture Placeholder 5", "role": "col_2_image_bottom", "required": False, "accepts": "image"},
            {"idx": 18, "name": "Picture Placeholder 5", "role": "col_3_image_bottom", "required": False, "accepts": "image"},
        ],
    },
    "12_Title and Content": {
        "templateId": "content_image_text_2col",
        "label": "Two Column Image and Text",
        "useWhen": "Two columns with images above and text labels below each",
        "compatibleWith": None,
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 14, "name": "Text Placeholder 3", "role": "subtitle", "required": False, "accepts": "text"},
            {"idx": 16, "name": "Picture Placeholder 5", "role": "left_image", "required": False, "accepts": "image"},
            {"idx": 12, "name": "Picture Placeholder 5", "role": "right_image", "required": False, "accepts": "image"},
            {"idx": 18, "name": "Text Placeholder 3", "role": "left_label", "required": False, "accepts": "text"},
            {"idx": 17, "name": "Text Placeholder 3", "role": "right_label", "required": False, "accepts": "text"},
        ],
    },
    "16_Title and Content": {
        "templateId": "content_with_caption",
        "label": "Content with Bottom Caption",
        "useWhen": "Standard content slide with a caption or subtext area at the bottom",
        "compatibleWith": "content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 13, "name": "Text Placeholder 4", "role": "subtitle", "required": False, "accepts": "text"},
            {"idx": 1, "name": "Content Placeholder 2", "role": "body", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "1_Picture with Caption": {
        "templateId": "image_with_text",
        "label": "Image with Text Side-by-Side",
        "useWhen": "Left text area with title + right image/content area",
        "compatibleWith": "content_image_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 13, "name": "Content Placeholder 3", "role": "body", "required": False, "accepts": "bullets|text|richtext"},
            {"idx": 14, "name": "Content Placeholder 7", "role": "image_area", "required": False, "accepts": "image|richtext"},
        ],
    },
    "2_Picture with Caption": {
        "templateId": "image_with_text_stacked",
        "label": "Image with Stacked Text",
        "useWhen": "Left text with title, right side has two stacked content areas",
        "compatibleWith": "content_image_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 13, "name": "Content Placeholder 3", "role": "body", "required": False, "accepts": "bullets|text|richtext"},
            {"idx": 14, "name": "Content Placeholder 7", "role": "image_area_top", "required": False, "accepts": "image|richtext"},
            {"idx": 15, "name": "Content Placeholder 7", "role": "image_area_bottom", "required": False, "accepts": "image|richtext"},
        ],
    },
    "7_Two Content": {
        "templateId": "content_two_with_header",
        "label": "Two Column with Header Bar",
        "useWhen": "Two equal content columns with a subtitle/header bar spanning the top",
        "compatibleWith": "two_content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 18, "name": "Text Placeholder 4", "role": "subtitle", "required": False, "accepts": "text"},
            {"idx": 17, "name": "Content Placeholder 2", "role": "left_content", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 19, "name": "Content Placeholder 2", "role": "right_content", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "Comparison": {
        "templateId": "comparison",
        "label": "Comparison Layout",
        "useWhen": "Direct comparison of two items with labels and content areas",
        "compatibleWith": "comparison_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 1, "name": "Text Placeholder 2", "role": "left_subtitle", "required": False, "accepts": "text"},
            {"idx": 2, "name": "Content Placeholder 3", "role": "left_content", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 3, "name": "Text Placeholder 4", "role": "right_subtitle", "required": False, "accepts": "text"},
            {"idx": 4, "name": "Content Placeholder 5", "role": "right_content", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "13_Title and Content": {
        "templateId": "content_stacked_compact",
        "label": "Stacked Compact",
        "useWhen": "Two stacked content areas in a compact layout",
        "compatibleWith": "content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 1, "name": "Content Placeholder 2", "role": "body_top", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 10, "name": "Content Placeholder 2", "role": "body_bottom", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "17_Title and Content": {
        "templateId": "content_stacked_compact_alt",
        "label": "Stacked Compact Alternate",
        "useWhen": "Two stacked content areas, same structure as 13 variant",
        "compatibleWith": "content_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 1, "name": "Content Placeholder 2", "role": "body_top", "required": True, "accepts": "bullets|text|richtext"},
            {"idx": 10, "name": "Content Placeholder 2", "role": "body_bottom", "required": True, "accepts": "bullets|text|richtext"},
        ],
    },
    "End": {
        "templateId": "closing_slide",
        "label": "Closing Slide",
        "useWhen": "Final slide of any deck, thank-you or Q&A",
        "compatibleWith": "closing_slide",
        "slots": [
            {"idx": 0, "name": "Title 1", "role": "slide_title", "required": True, "accepts": "text"},
            {"idx": 1, "name": "Subtitle 2", "role": "subtitle", "required": False, "accepts": "text"},
        ],
    },
    "DEFAULT": {
        "templateId": "default",
        "label": "Default Layout",
        "useWhen": "Fallback default layout with no defined placeholders",
        "compatibleWith": None,
        "slots": [],
    },
}


def _build_ext_element(meta: Dict[str, Any]) -> etree._Element:
    ext = etree.Element(f"{{{_P_NS}}}ext")
    ext.set("uri", CUSTOM_URI)

    nsmap = {"meta": CUSTOM_NS}
    root = etree.SubElement(ext, f"{{{CUSTOM_NS}}}layoutMeta", nsmap=nsmap)

    etree.SubElement(root, f"{{{CUSTOM_NS}}}templateId").text = meta["templateId"]
    etree.SubElement(root, f"{{{CUSTOM_NS}}}label").text = meta["label"]
    etree.SubElement(root, f"{{{CUSTOM_NS}}}useWhen").text = meta["useWhen"]

    compat = meta.get("compatibleWith")
    if compat:
        etree.SubElement(root, f"{{{CUSTOM_NS}}}compatibleWith").text = compat

    slots_el = etree.SubElement(root, f"{{{CUSTOM_NS}}}slots")
    for s in meta.get("slots", []):
        slot = etree.SubElement(slots_el, f"{{{CUSTOM_NS}}}slot")
        slot.set("idx", str(s["idx"]))
        slot.set("name", s["name"])
        slot.set("role", s["role"])
        slot.set("required", "true" if s.get("required") else "false")
        slot.set("accepts", s.get("accepts", "text"))

    return ext


def embed_metadata(input_path: str, output_path: str) -> None:
    prs = Presentation(input_path)

    matched = 0
    skipped = 0

    for layout in prs.slide_layouts:
        meta = LAYOUT_METADATA.get(layout.name)
        if not meta:
            skipped += 1
            print(f"  [skip] No metadata for layout: '{layout.name}'")
            continue

        el = layout._element

        extLst = el.find(f"{{{_P_NS}}}extLst")
        if extLst is not None:
            for ext in extLst.findall(f"{{{_P_NS}}}ext"):
                if ext.get("uri") == CUSTOM_URI:
                    extLst.remove(ext)
        else:
            extLst = etree.SubElement(el, f"{{{_P_NS}}}extLst")

        extLst.append(_build_ext_element(meta))
        matched += 1
        print(f"  [ok]   '{layout.name}' -> templateId={meta['templateId']}")

    prs.save(output_path)
    print(f"\nEmbedded {matched} layouts, skipped {skipped}")
    print(f"Saved: {output_path}")


def main() -> None:
    src = sys.argv[1] if len(sys.argv) > 1 else str(TEMPLATES_DIR / "template.pptx")
    dst = sys.argv[2] if len(sys.argv) > 2 else str(TEMPLATES_DIR / "template_tagged.pptx")

    if not Path(src).exists():
        print(f"Error: input file not found: {src}")
        sys.exit(1)

    print(f"Reading: {src}")
    embed_metadata(src, dst)


if __name__ == "__main__":
    main()
