"""
inspect_template.py
===================
Reads <p:ext> layout metadata from a template .pptx and outputs a JSON manifest.

If the template has embedded metadata (from embed_metadata.py), this script reads
it. Otherwise, it falls back to scanning placeholder structure heuristically.

Usage:
    python scripts/inspect_template.py [template.pptx]
    python scripts/inspect_template.py scripts/templates/template_tagged.pptx
"""

from __future__ import annotations

import json
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

_SCRIPT_DIR = Path(__file__).resolve().parent
TEMPLATES_DIR = _SCRIPT_DIR / "templates"

CUSTOM_NS = "https://beteek.com/pptx-layout-meta"
CUSTOM_URI = "{b9e7a3d1-4c5f-4a8b-b2d6-e1f3a7c9d0e2}"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

_PH_TYPE_LABELS = {
    1: "TITLE", 2: "BODY", 3: "CENTER_TITLE", 4: "SUBTITLE",
    7: "OBJECT", 18: "PICTURE", 13: "SLIDE_NUMBER",
    16: "DATE", 15: "FOOTER",
}

_META_PH_TYPES = {13, 16, 15}


def _read_layout_metadata(layout) -> Optional[Dict[str, Any]]:
    extLst = layout._element.find(f"{{{_P_NS}}}extLst")
    if extLst is None:
        return None

    for ext in extLst:
        if ext.get("uri") != CUSTOM_URI:
            continue

        tmpl = ext.find(f"{{{CUSTOM_NS}}}layoutMeta")
        if tmpl is None:
            continue

        def _text(tag: str) -> str:
            el = tmpl.find(f"{{{CUSTOM_NS}}}{tag}")
            return el.text.strip() if el is not None and el.text else ""

        compat_el = tmpl.find(f"{{{CUSTOM_NS}}}compatibleWith")
        compat = compat_el.text.strip() if compat_el is not None and compat_el.text else None

        slots = []
        for slot_el in tmpl.findall(f".//{{{CUSTOM_NS}}}slot"):
            slots.append({
                "idx": int(slot_el.get("idx", "0")),
                "name": slot_el.get("name", ""),
                "role": slot_el.get("role", ""),
                "required": slot_el.get("required", "false") == "true",
                "accepts": slot_el.get("accepts", "text"),
            })

        return {
            "templateId": _text("templateId"),
            "label": _text("label"),
            "useWhen": _text("useWhen"),
            "compatibleWith": compat,
            "slots": slots,
        }

    return None


def _describe_placeholders(layout) -> List[Dict[str, Any]]:
    result = []
    for ph in layout.placeholders:
        ptype_num = int(ph.placeholder_format.type)
        ptype = _PH_TYPE_LABELS.get(ptype_num, str(ptype_num))
        result.append({
            "idx": ph.placeholder_format.idx,
            "name": ph.name,
            "type": ptype,
            "isMeta": ptype_num in _META_PH_TYPES,
        })
    return result


def inspect_template(pptx_path: str) -> Dict[str, Any]:
    prs = Presentation(pptx_path)
    layouts = []

    for i, layout in enumerate(prs.slide_layouts):
        meta = _read_layout_metadata(layout)
        entry: Dict[str, Any] = {
            "layoutIndex": i,
            "layoutName": layout.name,
        }

        if meta:
            entry["hasMetadata"] = True
            entry["templateId"] = meta["templateId"]
            entry["label"] = meta["label"]
            entry["useWhen"] = meta["useWhen"]
            if meta["compatibleWith"]:
                entry["compatibleWith"] = meta["compatibleWith"]
            entry["slots"] = meta["slots"]
        else:
            entry["hasMetadata"] = False
            entry["placeholders"] = _describe_placeholders(layout)

        layouts.append(entry)

    content_slots = []
    for lay in layouts:
        if lay.get("hasMetadata"):
            slots = [s for s in lay.get("slots", []) if s["role"] != "slide_title"]
            entry = {
                "templateId": lay["templateId"],
                "layoutName": lay["layoutName"],
                "label": lay["label"],
                "contentSlotCount": len(slots),
                "slotRoles": [s["role"] for s in slots],
            }
            if lay.get("compatibleWith"):
                entry["compatibleWith"] = lay["compatibleWith"]
            content_slots.append(entry)

    return {
        "sourceFile": str(pptx_path),
        "totalLayouts": len(layouts),
        "layouts": layouts,
        "summary": {
            "withMetadata": sum(1 for l in layouts if l.get("hasMetadata")),
            "withoutMetadata": sum(1 for l in layouts if not l.get("hasMetadata")),
            "availableTemplateIds": [
                l["templateId"] for l in layouts if l.get("hasMetadata")
            ],
        },
    }


def main() -> None:
    default_path = str(TEMPLATES_DIR / "template_tagged.pptx")
    fallback_path = str(TEMPLATES_DIR / "template.pptx")

    path = sys.argv[1] if len(sys.argv) > 1 else default_path
    if not Path(path).exists():
        if Path(fallback_path).exists():
            print(f"Note: {path} not found, falling back to {fallback_path}", file=sys.stderr)
            path = fallback_path
        else:
            print(f"Error: no template found at {path} or {fallback_path}", file=sys.stderr)
            sys.exit(1)

    result = inspect_template(path)
    print(json.dumps(result, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
