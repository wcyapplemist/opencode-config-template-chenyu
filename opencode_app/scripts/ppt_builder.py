"""
ppt_builder.py
==============
PPT engine using template.pptx Slide Master layouts with proper placeholders.

Loads the template, adds new slides from named layouts (resolved by name, not
index, so layout reordering does not break it), fills placeholders by type
(TITLE, SUBTITLE, OBJECT), and saves the result.

Layouts are matched by name via ``_LAYOUT_NAME_MAP``; ``template.config.json``
may override the layout name for ``title_slide`` / ``content_slide``.

If a ``template_tagged.pptx`` (with embedded ``<p:ext>`` layout metadata) is
present, the engine reads metadata from each layout to build the mapping
dynamically. This makes the template self-describing — swap templates without
changing code. Falls back to the hardcoded ``_LAYOUT_NAME_MAP`` for plain
templates.

Usage:
    from ppt_builder import generate_ppt_from_data, DEFAULT_OUTPUT_DIR

    result = generate_ppt_from_data(
        slide_data_list,
        output_path=str(DEFAULT_OUTPUT_DIR / "report.pptx"),
    )
"""

from __future__ import annotations

import json
import logging
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)

_SCRIPT_DIR = Path(__file__).resolve().parent
TEMPLATES_DIR = _SCRIPT_DIR / "templates"
DEFAULT_OUTPUT_DIR = _SCRIPT_DIR.parent / "output"

_TEMPLATE_FILE = TEMPLATES_DIR / "template.pptx"
_TEMPLATE_TAGGED = TEMPLATES_DIR / "template_tagged.pptx"

_CUSTOM_META_NS = "https://beteek.com/pptx-layout-meta"
_CUSTOM_META_URI = "{b9e7a3d1-4c5f-4a8b-b2d6-e1f3a7c9d0e2}"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

_TITLE_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
_SUBTITLE_TYPE = PP_PLACEHOLDER.SUBTITLE
_BODY_TYPE = PP_PLACEHOLDER.BODY
_OBJECT_TYPE = PP_PLACEHOLDER.OBJECT

_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

_LAYOUT_NAME_MAP: Dict[str, List[str]] = {
    "title_slide": ["Title Slide"],
    "closing_slide": ["End"],
    "section_header_slide": ["Section Header"],
    "content_slide": ["Title and Content"],
    "two_content_slide": ["7_Two Content"],
    "comparison_slide": ["Comparison"],
    "content_image_slide": ["Picture with Caption"],
}

_LAYOUTS_WITH_SUBTITLE = {
    "title_slide", "closing_slide",
}
_LAYOUTS_WITH_BODY = {
    "content_slide", "content_image_slide",
}
_LAYOUTS_WITH_TWO_BODIES = {
    "two_content_slide", "comparison_slide",
}


def _resolve_template(template_arg: Optional[str] = None) -> Path:
    if template_arg and template_arg != "auto":
        p = Path(template_arg)
        if p.exists():
            return p
    if _TEMPLATE_TAGGED.exists():
        logger.info("Using tagged template: %s", _TEMPLATE_TAGGED.name)
        return _TEMPLATE_TAGGED
    logger.info("Using plain template: %s", _TEMPLATE_FILE.name)
    return _TEMPLATE_FILE


def _read_layout_metadata(prs: Presentation) -> Dict[str, Any]:
    result: Dict[str, Any] = {
        "by_template_id": {},
        "by_compatible_with": {},
        "layout_slots": {},
    }
    _ns = _CUSTOM_META_NS
    _ns_tag = "{" + _ns + "}"
    for layout in prs.slide_layouts:
        el = layout._element
        extLst = el.find(f"{{{_P_NS}}}extLst")
        if extLst is None:
            continue
        for ext in extLst:
            if ext.get("uri") != _CUSTOM_META_URI:
                continue
            tmpl = ext.find(f"{_ns_tag}layoutMeta")
            if tmpl is None:
                continue

            def _text(tag: str) -> str:
                e = tmpl.find(f"{_ns_tag}{tag}")
                return e.text.strip() if e is not None and e.text else ""

            compat_el = tmpl.find(f"{_ns_tag}compatibleWith")
            compat = compat_el.text.strip() if compat_el is not None and compat_el.text else None

            template_id = _text("templateId")

            slots = []
            for slot_el in tmpl.findall(f".//{_ns_tag}slot"):
                slots.append({
                    "idx": int(slot_el.get("idx", "0")),
                    "name": slot_el.get("name", ""),
                    "role": slot_el.get("role", ""),
                    "required": slot_el.get("required", "false") == "true",
                    "accepts": slot_el.get("accepts", "text"),
                })

            if template_id:
                result["by_template_id"][template_id] = layout
                result["layout_slots"][template_id] = slots
            if compat:
                result["by_compatible_with"].setdefault(compat, []).append(layout)
            break
    return result


def _normalize_layout_name(name: str) -> str:
    return re.sub(r"^\d+_", "", name).strip().lower()


def _build_layout_index(prs: Presentation) -> Tuple[Dict[str, Any], Dict[str, Any]]:
    exact: Dict[str, Any] = {}
    normalized: Dict[str, Any] = {}
    for layout in prs.slide_layouts:
        nm = layout.name
        exact[nm.lower()] = layout
        norm_key = _normalize_layout_name(nm)
        normalized.setdefault(norm_key, layout)
    return exact, normalized


def _resolve_layout(
    candidate_names: List[str],
    exact: Dict[str, Any],
    normalized: Dict[str, Any],
) -> Optional[Any]:
    for cand in candidate_names:
        if cand.lower() in exact:
            return exact[cand.lower()]
    for cand in candidate_names:
        key = _normalize_layout_name(cand)
        if key in normalized:
            return normalized[key]
    return None


def _load_config() -> Dict[str, Any]:
    config_path = TEMPLATES_DIR / "template.config.json"
    if config_path.exists():
        with open(config_path, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def _remove_all_slides(prs: Presentation) -> int:
    count = len(prs.slides)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].attrib.get(f"{{{_REL_NS}}}id")
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    return count


def _find_placeholder(slide: Any, ph_type: Any) -> Optional[Any]:
    for ph in slide.placeholders:
        if ph.placeholder_format.type == ph_type:
            return ph
    return None


def _find_placeholders(slide: Any, ph_type: Any) -> List[Any]:
    return [ph for ph in slide.placeholders if ph.placeholder_format.type == ph_type]


def _find_title_placeholder(slide: Any) -> Optional[Any]:
    for ph_type in _TITLE_TYPES:
        ph = _find_placeholder(slide, ph_type)
        if ph:
            return ph
    return None


def _find_body_placeholder(slide: Any) -> Optional[Any]:
    ph = _find_placeholder(slide, _BODY_TYPE)
    if ph:
        return ph
    objects = _find_placeholders(slide, _OBJECT_TYPE)
    return objects[0] if objects else None


def _set_text(shape: Any, text: str) -> bool:
    if not shape or not shape.has_text_frame:
        return False
    try:
        tf = shape.text_frame
        tf.clear()
        tf.paragraphs[0].text = text
        return True
    except Exception as exc:
        logger.warning("Failed to set text: %s", exc)
        return False


def _parse_line(line: str) -> Tuple[str, str]:
    clean = re.sub(r"\*\*", "", line.strip())
    if not clean:
        return ("", "")
    for sep in [" \u2014 ", " - ", ": "]:
        if sep in clean:
            parts = clean.split(sep, 1)
            return (parts[0].strip(), parts[1].strip())
    return (clean, "")


def _set_notes(slide: Any, notes_text: str) -> bool:
    text = (notes_text or "").strip()
    if not text:
        return False
    try:
        slide.notes_slide.notes_text_frame.text = text
        return True
    except Exception as exc:
        logger.warning("Failed to set notes: %s", exc)
        return False


def _set_body_text(shape: Any, text: str) -> bool:
    if not shape or not shape.has_text_frame:
        return False
    try:
        lines = [l for l in text.split("\n") if l.strip()]
        tf = shape.text_frame
        tf.clear()

        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            title_part, desc_part = _parse_line(line)

            if title_part:
                run = p.add_run()
                run.text = title_part
                run.font.bold = True
                run.font.size = Pt(14)
            if desc_part:
                run = p.add_run()
                run.text = f" \u2014 {desc_part}" if title_part else desc_part
                run.font.size = Pt(12)

        return True
    except Exception as exc:
        logger.warning("Failed to set body text: %s", exc)
        return False


def _find_placeholder_by_idx(slide: Any, idx: int) -> Optional[Any]:
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _fill_by_slot_roles(
    slide: Any,
    slide_data: Dict[str, Any],
    layout_slots: List[Dict[str, Any]],
) -> int:
    slot_values = slide_data.get("slots", {})
    filled = 0

    for slot_meta in layout_slots:
        role = slot_meta["role"]
        if role == "slide_title":
            continue
        idx = slot_meta["idx"]
        accepts = slot_meta.get("accepts", "text")

        content = slot_values.get(role, "")
        if not content:
            continue

        ph = _find_placeholder_by_idx(slide, idx)
        if ph is None:
            logger.warning("  Slot '%s' (idx=%d): placeholder not found", role, idx)
            continue

        is_body = "bullets" in accepts or "richtext" in accepts
        ok = _set_body_text(ph, content) if is_body else _set_text(ph, content)
        if ok:
            filled += 1
            logger.info("  Slot '%s': filled (%s)", role, "body" if is_body else "text")

    return filled


def _resolve_layout_metadata(
    slide_type: str,
    slide_data: Dict[str, Any],
    metadata: Dict[str, Any],
    config: Dict[str, Any],
    exact_idx: Dict[str, Any],
    norm_idx: Dict[str, Any],
) -> Optional[Any]:
    direct_id = slide_data.get("template_id") or slide_type
    if direct_id in metadata["by_template_id"]:
        return metadata["by_template_id"][direct_id]

    compat_list = metadata["by_compatible_with"].get(slide_type)
    if compat_list:
        if slide_type == "title_slide" and config.get("title_slide_layout"):
            name = config["title_slide_layout"]
            for lay in compat_list:
                if lay.name.lower() == name.lower():
                    return lay
            cand = _resolve_layout([name], exact_idx, norm_idx)
            if cand:
                return cand
        elif slide_type == "content_slide" and config.get("content_slide_layout"):
            name = config["content_slide_layout"]
            for lay in compat_list:
                if lay.name.lower() == name.lower():
                    return lay
        return compat_list[0]

    return None


def _get_layout_template_id(layout: Any, metadata: Dict[str, Any]) -> Optional[str]:
    for tid, lay in metadata["by_template_id"].items():
        if lay is layout:
            return tid
    return None


def generate_ppt_from_data(
    slide_data_list: List[Dict[str, Any]],
    template_path: Optional[str] = None,
    output_path: str = "output.pptx",
    prompt_text: str = "",
) -> str:
    template = _resolve_template(template_path)
    output = Path(output_path)

    if not template.exists():
        raise FileNotFoundError(f"Template not found: {template}")

    if not output.is_absolute():
        output = DEFAULT_OUTPUT_DIR / output
    output.parent.mkdir(parents=True, exist_ok=True)

    config = _load_config()

    logger.info("Loading template: %s", template.name)
    prs = Presentation(str(template))
    logger.info("Template: %d slides, %d layouts", len(prs.slides), len(prs.slide_layouts))

    removed = _remove_all_slides(prs)
    logger.info("Cleared %d example slides", removed)

    exact_idx, norm_idx = _build_layout_index(prs)
    metadata = _read_layout_metadata(prs)
    has_metadata = bool(metadata["by_template_id"])
    if has_metadata:
        logger.info(
            "Metadata: %d layouts with templateId, %d with compatibleWith",
            len(metadata["by_template_id"]),
            sum(len(v) for v in metadata["by_compatible_with"].values()),
        )

    for page_num, slide_data in enumerate(slide_data_list, start=1):
        slide_type = slide_data.get("slide_type", "")

        layout = None

        if has_metadata:
            layout = _resolve_layout_metadata(slide_type, slide_data, metadata, config, exact_idx, norm_idx)

        if layout is None:
            if slide_type == "title_slide" and config.get("title_slide_layout"):
                candidates = [config["title_slide_layout"]]
            elif slide_type == "content_slide" and config.get("content_slide_layout"):
                candidates = [config["content_slide_layout"]]
            else:
                candidates = _LAYOUT_NAME_MAP.get(slide_type)

            if not candidates:
                logger.warning("Page %d: unknown slide_type '%s', skipped", page_num, slide_type)
                continue

            layout = _resolve_layout(candidates, exact_idx, norm_idx)

        if layout is None:
            logger.warning(
                "Page %d: no layout matched for slide_type '%s', skipped",
                page_num, slide_type,
            )
            continue

        layout_idx = prs.slide_layouts.index(layout)
        try:
            slide = prs.slides.add_slide(layout)
            logger.info("Page %d: added slide from layout[%d] '%s'", page_num, layout_idx, layout.name)

            title_text = slide_data.get("title", "")
            use_slots = "slots" in slide_data

            if use_slots:
                if title_text:
                    title_ph = _find_title_placeholder(slide)
                    if title_ph:
                        _set_text(title_ph, title_text)
                        logger.info("  Title: \"%s\"", title_text)

                layout_tid = _get_layout_template_id(layout, metadata)
                layout_slots_meta = metadata["layout_slots"].get(layout_tid, [])
                _fill_by_slot_roles(slide, slide_data, layout_slots_meta)
            else:
                if title_text:
                    title_ph = _find_title_placeholder(slide)
                    if title_ph:
                        _set_text(title_ph, title_text)
                        logger.info("  Title: \"%s\"", title_text)

                if slide_type in _LAYOUTS_WITH_SUBTITLE:
                    subtitle_text = slide_data.get("subtitle", "")
                    if subtitle_text:
                        sub_ph = _find_placeholder(slide, _SUBTITLE_TYPE)
                        if sub_ph:
                            _set_text(sub_ph, subtitle_text)
                            logger.info("  Subtitle: \"%s\"", subtitle_text[:50])

                if slide_type in _LAYOUTS_WITH_BODY:
                    body_text = slide_data.get("body", "")
                    if body_text:
                        body_ph = _find_body_placeholder(slide)
                        if body_ph:
                            _set_body_text(body_ph, body_text)
                            logger.info("  Body: %d lines", len([l for l in body_text.split("\n") if l.strip()]))

                if slide_type in _LAYOUTS_WITH_TWO_BODIES:
                    body_left = slide_data.get("body_left", "")
                    body_right = slide_data.get("body_right", "")
                    objects = _find_placeholders(slide, _OBJECT_TYPE)
                    if len(objects) >= 2:
                        if body_left:
                            _set_body_text(objects[0], body_left)
                            logger.info("  Body-left: %d lines", len([l for l in body_left.split("\n") if l.strip()]))
                        if body_right:
                            _set_body_text(objects[1], body_right)
                            logger.info("  Body-right: %d lines", len([l for l in body_right.split("\n") if l.strip()]))
                    elif len(objects) == 1 and (body_left or body_right):
                        _set_body_text(objects[0], body_left or body_right)

            notes_text = slide_data.get("notes", "")
            if _set_notes(slide, notes_text):
                logger.info("  Notes: %d chars", len(notes_text))

        except Exception as exc:
            logger.error("Page %d failed: %s", page_num, exc)

    prs.save(str(output))
    logger.info("Saved: %s (%d slides)", output.resolve(), len(prs.slides))
    return str(output.resolve())


def main() -> None:
    mock: List[Dict[str, Any]] = [
        {
            "slide_type": "title_slide",
            "title": "AI Empowering Finance",
            "subtitle": "2026 Q1",
            "notes": (
                "KEY MESSAGE: Open with energy — set the stakes in one line.\n"
                "\"Hold the slide for two seconds before you speak.\"\n"
                "\"Good [morning/afternoon], I'm [Name]. Today I want to show you how AI is already transforming finance — not in theory, but in the numbers.\"\n"
                "Pause. Let the tagline land.\n"
                "\"We'll walk through where it delivers the clearest ROI today.\"\n"
                "TRANSITION: \"Let me start with the core scenarios.\"\n"
                "COACHING: Eye contact, confident. Do not read the slide. Be ready for: \"Is this hype or real?\" — lead with the 80 percent figure."
            ),
        },
        {
            "slide_type": "content_slide",
            "title": "Core AI Scenarios",
            "body": (
                "**Automated Reporting** \u2014 RPA tools auto-generate monthly reports, cutting manual effort by 80%\n"
                "**Smart Reconciliation** \u2014 AI matches bank transactions at 99.5% accuracy\n"
                "**Fraud Detection** \u2014 Real-time anomaly detection with automated alerts\n"
                "**Tax Optimization** \u2014 ML identifies savings opportunities across tax structures"
            ),
            "notes": (
                "KEY MESSAGE: Four high-impact scenarios where AI already delivers measurable ROI.\n"
                "\"Let's make this concrete. These aren't edge cases — this is everyday finance.\"\n"
                "\"Automated reporting alone removes eighty percent of the manual effort behind every monthly close.\"\n"
                "Pause. Let the number land.\n"
                "\"Smart reconciliation now matches transactions at ninety-nine-point-five percent accuracy, and fraud detection flags anomalies in real time.\"\n"
                "\"Ask your CFO: how much would one missed discrepancy cost?\"\n"
                "TRANSITION: \"Here is how we roll this out.\"\n"
                "COACHING: Matter-of-fact tone, don't over-sell. Be ready for: \"What about false positives?\" — answer: tuned thresholds, human-in-the-loop review."
            ),
        },
        {
            "slide_type": "content_slide",
            "title": "Roadmap",
            "body": (
                "**Phase 1: Pilot** \u2014 Deploy in 2 business units by Q2\n"
                "**Phase 2: Scale** \u2014 Expand to all departments by Q4\n"
                "**Phase 3: Full Deployment** \u2014 Organization-wide adoption by 2027"
            ),
            "notes": (
                "KEY MESSAGE: A phased, low-risk rollout — pilot, scale, then full adoption.\n"
                "\"We don't boil the ocean. We pilot in two units first, prove the numbers, then scale.\"\n"
                "\"By Q4 every department is on board, and full organisation-wide adoption lands in 2027.\"\n"
                "Walk the three phases left to right.\n"
                "TRANSITION: Open for questions.\n"
                "COACHING: Keep it tight, end with confidence. Be ready for: \"What could delay Phase 2?\" — answer: only change-management, never the technology."
            ),
        },
    ]

    print("Test: template.pptx (placeholder-based)")
    result = generate_ppt_from_data(
        mock,
        output_path=str(DEFAULT_OUTPUT_DIR / "test_template.pptx"),
    )
    print(f"Output: {result}")


if __name__ == "__main__":
    main()
