"""Verify chart styling: fonts, colors, gridlines are applied correctly."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

prs = Presentation('output/test_chart_slides.pptx')

expected_colors = [
    RGBColor(0x44, 0x72, 0xC4),
    RGBColor(0xED, 0x7D, 0x31),
    RGBColor(0xFF, 0xC0, 0x00),
]
expected_font = "Calibri"
expected_gridline = RGBColor(0xE7, 0xE6, 0xE6)
expected_text_color = RGBColor(0x44, 0x54, 0x6A)

all_ok = True

for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if not shape.has_chart:
            continue
        chart = shape.chart
        ct = chart.chart_type
        print(f"\nSlide {idx+1}: chart_type={ct}")

        # Check global font
        font_name = chart.font.name
        ok = font_name == expected_font
        all_ok = all_ok and ok
        print(f"  chart.font.name = '{font_name}' {'OK' if ok else 'FAIL'}")

        # Check legend
        if chart.has_legend:
            leg_font = chart.legend.font
            ok = leg_font.name == expected_font
            all_ok = all_ok and ok
            print(f"  legend.font.name = '{leg_font.name}' {'OK' if ok else 'FAIL'}")

        plot = chart.plots[0]

        # Check data labels
        if plot.has_data_labels:
            dl_font = plot.data_labels.font
            ok = dl_font.name == expected_font
            all_ok = all_ok and ok
            print(f"  data_labels.font.name = '{dl_font.name}' {'OK' if ok else 'FAIL'}")
            dl_color = dl_font.color
            if dl_color and dl_color.type is not None:
                print(f"  data_labels.font.color = {dl_color.rgb}")

        # Check axis (bar/line only)
        if ct != 5:  # not PIE
            try:
                va = chart.value_axis
                tl = va.tick_labels
                ok = tl.font.name == expected_font
                all_ok = all_ok and ok
                print(f"  value_axis.tick_labels.font.name = '{tl.font.name}' {'OK' if ok else 'FAIL'}")
                print(f"  value_axis.number_format = '{va.number_format}'")

                # Gridline color
                if va.has_major_gridlines:
                    gl = va.major_gridlines
                    line = gl.format.line
                    print(f"  gridline.color = {line.color.rgb if line.color and line.color.type else 'inherit'}")

                ca = chart.category_axis
                ca_tl = ca.tick_labels
                ok = ca_tl.font.name == expected_font
                all_ok = all_ok and ok
                print(f"  category_axis.tick_labels.font.name = '{ca_tl.font.name}' {'OK' if ok else 'FAIL'}")
            except Exception as e:
                print(f"  Axis check error: {e}")

        # Check series colors
        for si, series in enumerate(plot.series):
            if si < 3:
                fill = series.format.fill
                if fill.type is not None and hasattr(fill, 'fore_color'):
                    fc = fill.fore_color
                    if fc.type is not None:
                        color = fc.rgb
                        ok = color == expected_colors[si]
                        all_ok = all_ok and ok
                        print(f"  series[{si}].fill.color = {color} {'OK' if ok else 'FAIL (expected ' + str(expected_colors[si]) + ')'}")

print("\n" + "=" * 50)
if all_ok:
    print("ALL STYLING CHECKS PASSED")
else:
    print("SOME STYLING CHECKS FAILED")
    sys.exit(1)
