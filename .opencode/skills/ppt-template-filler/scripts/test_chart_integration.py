"""
Integration test: chart_slide types (bar, pie, line) + backward compatibility.

Tests that generate_ppt_from_data() correctly produces:
1. Bar chart slide (COLUMN_CLUSTERED)
2. Pie chart slide (PIE with percentage labels)
3. Line chart slide (LINE_MARKERS, multi-series)
4. Mixed deck: title + content + chart + closing (backward compat)
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
sys.path.insert(0, '.')
from ppt_builder import generate_ppt_from_data, DEFAULT_OUTPUT_DIR
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

OUTPUT = str(DEFAULT_OUTPUT_DIR / "test_chart_slides.pptx")

slide_data = [
    {
        "slide_type": "title_slide",
        "title": "Digital Technology in Construction",
        "subtitle": "Chart Integration Test",
    },
    {
        "slide_type": "content_slide",
        "title": "Industry Pain Points",
        "body": (
            "**Low Profit Margins** \u2014 Average 5% across the industry\n"
            "**Schedule Overruns** \u2014 70% of projects exceed deadlines\n"
            "**Safety Incidents** \u2014 High accident rates on site"
        ),
    },
    {
        "slide_type": "chart_slide",
        "title": "Global Construction Tech Market (USD Billion)",
        "chart_type": "bar",
        "categories": ["2020", "2021", "2022", "2023", "2024", "2025", "2026"],
        "series": [
            {"name": "Market Size", "values": [8.5, 11.2, 14.8, 19.5, 25.1, 31.7, 39.4]},
        ],
        "chart_options": {
            "legend_position": "bottom",
            "show_data_labels": True,
            "y_axis_min": 0,
            "y_axis_max": 45,
        },
        "notes": "KEY MESSAGE: Market growing from 8.5B to 39.4B.",
    },
    {
        "slide_type": "chart_slide",
        "title": "Technology Adoption Rate by Category",
        "chart_type": "pie",
        "categories": ["BIM", "IoT", "Drones", "AI & ML", "Robotics", "Cloud"],
        "series": [
            {"name": "Adoption %", "values": [68, 45, 52, 28, 15, 72]},
        ],
        "chart_options": {
            "legend_position": "right",
            "show_data_labels": True,
        },
        "notes": "KEY MESSAGE: Cloud and BIM lead adoption.",
    },
    {
        "slide_type": "chart_slide",
        "title": "Project Performance Improvement (%)",
        "chart_type": "line_markers",
        "categories": ["2019", "2020", "2021", "2022", "2023", "2024", "2025"],
        "series": [
            {"name": "Cost Savings",       "values": [5, 8, 12, 16, 20, 25, 30]},
            {"name": "Schedule Reduction", "values": [3, 6, 10, 14, 19, 24, 28]},
            {"name": "Safety Improvement", "values": [2, 4, 8, 12, 18, 22, 27]},
        ],
        "chart_options": {
            "legend_position": "bottom",
            "show_data_labels": True,
            "y_axis_min": 0,
            "y_axis_max": 35,
        },
        "notes": "KEY MESSAGE: All metrics improving consistently.",
    },
    {
        "slide_type": "closing_slide",
        "title": "Thank You",
        "subtitle": "Questions & Discussion",
    },
]

print("=" * 60)
print("Running integration test: chart_slide types")
print("=" * 60)

result = generate_ppt_from_data(slide_data, output_path=OUTPUT)
print(f"\nGenerated: {result}")

print("\n" + "=" * 60)
print("Verification: reading back the output")
print("=" * 60)

prs = Presentation(result)
assert len(prs.slides) == 6, f"Expected 6 slides, got {len(prs.slides)}"

chart_checks = {
    2: ("Bar Chart", XL_CHART_TYPE.COLUMN_CLUSTERED),
    3: ("Pie Chart", XL_CHART_TYPE.PIE),
    4: ("Line Chart", XL_CHART_TYPE.LINE_MARKERS),
}

all_passed = True

for idx, slide in enumerate(prs.slides):
    title = ""
    for shape in slide.shapes:
        if shape.has_text_frame and shape == slide.shapes.title:
            title = slide.shapes.title.text
            break

    chart_found = False
    chart_type_val = None
    for shape in slide.shapes:
        if shape.has_chart:
            chart_found = True
            chart_type_val = shape.chart.chart_type
            break

    notes_ok = slide.has_notes_slide and bool(slide.notes_slide.notes_text_frame.text.strip())

    if idx in chart_checks:
        label, expected_type = chart_checks[idx]
        chart_ok = chart_found and chart_type_val == expected_type
        status = "PASS" if chart_ok else "FAIL"
        if not chart_ok:
            all_passed = False
        print(f"  Slide {idx+1} [{label}]: title='{title}', chart_type={chart_type_val}, notes={notes_ok} -> {status}")
    else:
        no_chart_ok = not chart_found
        status = "PASS" if no_chart_ok else "FAIL"
        if not no_chart_ok:
            all_passed = False
        print(f"  Slide {idx+1} [{title}]: no_chart={no_chart_ok}, notes={notes_ok} -> {status}")

print("\n" + "=" * 60)
if all_passed:
    print("ALL TESTS PASSED")
else:
    print("SOME TESTS FAILED")
    sys.exit(1)
print("=" * 60)
