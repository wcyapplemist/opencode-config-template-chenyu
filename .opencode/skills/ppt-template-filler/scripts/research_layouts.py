"""Research script: examine template Blank layouts in detail."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
sys.path.insert(0, '.')
from pptx import Presentation
from pptx.util import Inches, Emu, Pt

prs = Presentation('templates/template.pptx')

print('Slide dimensions:')
print(f'  Width:  {prs.slide_width} EMU = {prs.slide_width / 914400:.2f} in')
print(f'  Height: {prs.slide_height} EMU = {prs.slide_height / 914400:.2f} in')

# Examine all Blank-like layouts
blank_indices = [6, 9]  # 1_Blank and Blank
for idx in blank_indices:
    layout = prs.slide_layouts[idx]
    print(f'\n===== Layout [{idx}]: {layout.name} =====')
    print(f'  Placeholders: {len(layout.placeholders)}')
    for ph in layout.placeholders:
        print(f'    idx={ph.placeholder_format.idx} name="{ph.name}" type={ph.placeholder_format.type}')
        if ph.left is not None:
            print(f'      pos: left={ph.left/914400:.2f}in top={ph.top/914400:.2f}in w={ph.width/914400:.2f}in h={ph.height/914400:.2f}in')

# Check all shapes (non-placeholder shapes like background images/decorations)
for idx in [6, 9]:
    layout = prs.slide_layouts[idx]
    print(f'\n===== Layout [{idx}] {layout.name} - all shapes =====')
    for shape in layout.shapes:
        stype = shape.shape_type
        print(f'  Shape: "{shape.name}" type={stype}')
        if hasattr(shape, 'left') and shape.left is not None:
            print(f'    pos: left={shape.left/914400:.2f}in top={shape.top/914400:.2f}in w={shape.width/914400:.2f}in h={shape.height/914400:.2f}in')

# Also check the Title and Content layout for comparison (to understand title position)
print('\n===== Layout [10] Title and Content - title placeholder =====')
layout10 = prs.slide_layouts[10]
for ph in layout10.placeholders:
    if ph.placeholder_format.type is not None and 'TITLE' in str(ph.placeholder_format.type):
        print(f'  Title placeholder pos: left={ph.left/914400:.2f}in top={ph.top/914400:.2f}in w={ph.width/914400:.2f}in h={ph.height/914400:.2f}in')
