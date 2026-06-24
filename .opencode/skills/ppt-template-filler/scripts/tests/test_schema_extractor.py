"""Tests for schema_extractor.py (US-1.1).

Covers the Task-6 acceptance areas from PLAN-GIT-48:
  - bundled template extract -> validate passes
  - count invariant (63 layouts, every layout has a components array)
  - non-placeholder elements captured (layout 2: >=5 components vs 3 placeholders)
  - master components non-empty
  - polygon values in [0,1], exactly 4 points
  - font-cardinality rule (C1): non-text components have no font
  - group nesting recurses
  - empty/blank layout does not crash
  - synthetic edge cases (group, table, chart, zero-shape container)
  - negative test (non-PPTX input raises TemplateExtractionError)
  - second-template robustness (default Presentation())
"""
import os
from typing import Any, List, Optional

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER

from schema_extractor import (
    COMPONENT_TYPE_ENUM,
    PLACEHOLDER_TYPE_ENUM,
    TemplateExtractionError,
    _extract_components,
    _IdCounter,
    extract_schema,
    map_shape_type,
    normalize_polygon,
    validate_template_schema,
)


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------
@pytest.fixture
def schema(template_path):
    return extract_schema(template_path)


# ---------------------------------------------------------------------------
# (1) Bundled template extract -> validate passes
# ---------------------------------------------------------------------------
class TestExtractAndValidate:
    def test_validates_clean(self, schema):
        result = validate_template_schema(schema)
        assert result.is_valid, result.error_messages()

    def test_top_level_keys(self, schema):
        for key in (
            "template_metadata",
            "slide_master",
            "slide_layouts",
            "component_type_enum",
            "placeholder_type_enum",
        ):
            assert key in schema, f"missing top-level key {key}"

    def test_metadata_fields(self, schema):
        meta = schema["template_metadata"]
        assert meta["title"] and isinstance(meta["title"], str)
        assert meta["schema_version"]
        assert meta["generated_by"] == "opencode-pptx-subagent/schema_extractor"
        assert meta["generated_at"]
        dims = meta["slide_dimensions"]
        for k in ("width_emu", "height_emu", "width_inches", "height_inches", "aspect_ratio"):
            assert k in dims
        assert dims["width_emu"] > dims["height_emu"] > 0  # 16:9 landscape

    def test_enums_self_documenting(self, schema):
        assert schema["component_type_enum"] == COMPONENT_TYPE_ENUM
        assert schema["placeholder_type_enum"] == PLACEHOLDER_TYPE_ENUM


# ---------------------------------------------------------------------------
# (2) Count invariant + every layout has a components array
# ---------------------------------------------------------------------------
class TestCountInvariant:
    def test_layout_count_matches_pptx(self, template_path):
        prs = Presentation(template_path)
        schema = extract_schema(template_path)
        assert len(schema["slide_layouts"]) == len(prs.slide_layouts)

    def test_every_layout_has_components(self, schema):
        for L in schema["slide_layouts"]:
            assert isinstance(L["components"], list), L.get("layout_name")
            for key in ("layout_id", "layout_name", "layout_index"):
                assert key in L

    def test_layout_indices_sequential(self, schema):
        idxs = [L["layout_index"] for L in schema["slide_layouts"]]
        assert idxs == list(range(len(idxs)))


# ---------------------------------------------------------------------------
# (3) Non-placeholder elements captured (the core US-1.1 gap)
# ---------------------------------------------------------------------------
class TestNonPlaceholderCapture:
    def test_layout2_captures_non_placeholders(self, schema):
        """Layout 2 has 5 total shapes but only 3 placeholders -> 2 dropped
        before US-1.1. Now all 5 must appear as components."""
        layout2 = schema["slide_layouts"][2]
        assert len(layout2["components"]) >= 5

    def test_master_components_non_empty(self, schema):
        assert len(schema["slide_master"]["components"]) == 3

    def test_component_types_beyond_placeholder_present(self, schema):
        """Extraction must yield at least one non-placeholder type (image/shape)."""
        all_types = set()
        for L in schema["slide_layouts"]:
            for c in L["components"]:
                all_types.add(c["type"])
        for c in schema["slide_master"]["components"]:
            all_types.add(c["type"])
        assert "placeholder" in all_types
        assert all_types - {"placeholder"}, "no non-placeholder types captured"


# ---------------------------------------------------------------------------
# (4) Polygon correctness
# ---------------------------------------------------------------------------
class TestPolygon:
    def test_every_polygon_is_four_points_in_range(self, schema):
        for L in schema["slide_layouts"] + [schema["slide_master"]]:
            for c in L["components"]:
                poly = c["polygon"]
                assert len(poly) == 4, c["id"]
                for pt in poly:
                    assert set(pt.keys()) == {"x", "y"}, c["id"]
                    assert 0.0 <= pt["x"] <= 1.0, (c["id"], pt)
                    assert 0.0 <= pt["y"] <= 1.0, (c["id"], pt)

    def test_global_id_uniqueness(self, schema):
        ids = [c["id"] for L in schema["slide_layouts"] for c in L["components"]]
        ids += [c["id"] for c in schema["slide_master"]["components"]]
        assert len(ids) == len(set(ids)), "component ids not globally unique"

    def test_z_order_nonneg_int(self, schema):
        for L in schema["slide_layouts"]:
            for c in L["components"]:
                assert isinstance(c["z_order"], int) and c["z_order"] >= 0


# ---------------------------------------------------------------------------
# (5) Font-cardinality rule (architecture review C1)
# ---------------------------------------------------------------------------
TEXT_TYPES = {"textbox", "placeholder"}


class TestFontCardinality:
    def test_non_text_components_omit_font(self, schema):
        for L in schema["slide_layouts"] + [schema["slide_master"]]:
            for c in L["components"]:
                if c["type"] not in TEXT_TYPES:
                    assert "font" not in c, f"{c['id']} ({c['type']}) must not carry font"

    def test_text_components_carry_font_stub(self, schema):
        found_text = False
        for L in schema["slide_layouts"] + [schema["slide_master"]]:
            for c in L["components"]:
                if c["type"] in TEXT_TYPES:
                    found_text = True
                    assert "font" in c and c["font"] == {}, c["id"]
        assert found_text, "template produced no text components"


# ---------------------------------------------------------------------------
# (6) Validation catches violations
# ---------------------------------------------------------------------------
class TestValidationFinds:
    def test_rejects_font_on_non_text(self):
        bad = {
            "template_metadata": _ok_meta(),
            "slide_master": {"name": "M", "components": []},
            "slide_layouts": [{
                "layout_id": "x", "layout_name": "X", "layout_index": 0,
                "components": [_bad_component_with_font_on_image()],
            }],
            "component_type_enum": COMPONENT_TYPE_ENUM,
            "placeholder_type_enum": PLACEHOLDER_TYPE_ENUM,
        }
        r = validate_template_schema(bad)
        assert not r.is_valid
        assert any("must not carry a 'font'" in e.reason for e in r.errors)

    def test_rejects_polygon_out_of_range(self):
        comp = _ok_component()
        comp["polygon"][0]["x"] = 1.5
        r = validate_template_schema(_schema_with(comp))
        assert not r.is_valid
        assert any("out of [0,1]" in e.reason for e in r.errors)

    def test_rejects_bad_type_enum(self):
        comp = _ok_component()
        comp["type"] = "bogus"
        r = validate_template_schema(_schema_with(comp))
        assert not r.is_valid
        assert any("not in component_type_enum" in e.reason for e in r.errors)

    def test_rejects_non_object(self):
        r = validate_template_schema("not a dict")
        assert not r.is_valid


def _ok_meta():
    return {
        "title": "T", "schema_version": "1.0.0",
        "generated_by": "x", "generated_at": "2026-01-01T00:00:00Z",
        "slide_dimensions": {
            "width_emu": 9144000, "height_emu": 5143500,
            "width_inches": 10.0, "height_inches": 5.625, "aspect_ratio": "16:9",
        },
    }


def _ok_component():
    return {
        "id": "comp_001", "type": "textbox", "name": "N",
        "polygon": [{"x": 0.1, "y": 0.1}] * 4, "z_order": 0,
        "placeholder_type": None, "font": {}, "runs": [],
        "content_template": "{{content}}",
    }


def _bad_component_with_font_on_image():
    c = _ok_component()
    c["id"] = "comp_002"
    c["type"] = "image"
    c["font"] = {"family": "Arial"}  # image must not carry font
    return c


def _schema_with(comp):
    return {
        "template_metadata": _ok_meta(),
        "slide_master": {"name": "M", "components": []},
        "slide_layouts": [{
            "layout_id": "x", "layout_name": "X", "layout_index": 0,
            "components": [comp],
        }],
        "component_type_enum": COMPONENT_TYPE_ENUM,
        "placeholder_type_enum": PLACEHOLDER_TYPE_ENUM,
    }


# ---------------------------------------------------------------------------
# (7) Synthetic edge cases — unit-test mapper + recursion with fake shapes
# ---------------------------------------------------------------------------
class FakePlaceholderFormat:
    def __init__(self, ptype: Any) -> None:
        self.type = ptype


class FakeShape:
    """Minimal stand-in for a python-pptx shape, for pure-logic unit tests."""

    def __init__(
        self,
        name: str,
        st: Any = None,
        *,
        is_placeholder: bool = False,
        has_table: bool = False,
        has_chart: bool = False,
        left: int = 0,
        top: int = 0,
        width: int = 0,
        height: int = 0,
        children: Optional[List["FakeShape"]] = None,
        ph_type: Any = None,
    ) -> None:
        self.name = name
        self.shape_type = st
        self.is_placeholder = is_placeholder
        self.has_table = has_table
        self.has_chart = has_chart
        self.left, self.top, self.width, self.height = left, top, width, height
        self._children = children or []
        self.placeholder_format = FakePlaceholderFormat(ph_type) if is_placeholder else None

    @property
    def shapes(self):
        return self._children


SLIDE_W = 9144000  # EMU (10in)
SLIDE_H = 5143500  # EMU


class TestTypeMapper:
    def test_table_detected_before_shape_type(self):
        s = FakeShape("t", st=None, has_table=True)
        assert map_shape_type(s) == "table"

    def test_chart_detected(self):
        s = FakeShape("c", st=None, has_chart=True)
        assert map_shape_type(s) == "chart"

    def test_placeholder_takes_precedence(self):
        s = FakeShape("p", st=MSO_SHAPE_TYPE.PICTURE, is_placeholder=True, ph_type=PP_PLACEHOLDER.TITLE)
        assert map_shape_type(s) == "placeholder"

    def test_picture_image(self):
        s = FakeShape("pic", st=MSO_SHAPE_TYPE.PICTURE)
        assert map_shape_type(s) == "image"

    def test_group(self):
        s = FakeShape("g", st=MSO_SHAPE_TYPE.GROUP)
        assert map_shape_type(s) == "group"

    def test_textbox(self):
        s = FakeShape("tb", st=MSO_SHAPE_TYPE.TEXT_BOX)
        assert map_shape_type(s) == "textbox"

    def test_media_video(self):
        s = FakeShape("m", st=MSO_SHAPE_TYPE.MEDIA)
        assert map_shape_type(s) == "video"

    def test_smartart(self):
        s = FakeShape("sm", st=MSO_SHAPE_TYPE.IGX_GRAPHIC)
        assert map_shape_type(s) == "smartart"

    def test_unknown_degrades_to_shape(self):
        s = FakeShape("u", st=None)
        assert map_shape_type(s) == "shape"


class TestPolygonNormalizer:
    def test_four_points_tl_tr_br_bl(self):
        s = FakeShape("s", left=0, top=0, width=SLIDE_W, height=SLIDE_H)
        poly = normalize_polygon(s, SLIDE_W, SLIDE_H)
        assert poly[0] == {"x": 0.0, "y": 0.0}      # TL
        assert poly[1] == {"x": 1.0, "y": 0.0}      # TR
        assert poly[2] == {"x": 1.0, "y": 1.0}      # BR
        assert poly[3] == {"x": 0.0, "y": 1.0}      # BL

    def test_clamps_overflow(self):
        # shape extends beyond slide -> clamped to 1.0
        s = FakeShape("s", left=SLIDE_W // 2, top=SLIDE_H // 2,
                      width=SLIDE_W, height=SLIDE_H)
        poly = normalize_polygon(s, SLIDE_W, SLIDE_H)
        assert all(0.0 <= pt["x"] <= 1.0 for pt in poly)
        assert all(0.0 <= pt["y"] <= 1.0 for pt in poly)

    def test_zero_dims_safe(self):
        s = FakeShape("s")
        poly = normalize_polygon(s, 0, 0)
        assert poly == [{"x": 0.0, "y": 0.0}] * 4


class TestGroupRecursion:
    def test_group_children_flattened(self):
        child1 = FakeShape("c1", st=MSO_SHAPE_TYPE.TEXT_BOX,
                           left=0, top=0, width=1000, height=1000)
        child2 = FakeShape("c2", st=MSO_SHAPE_TYPE.PICTURE,
                           left=2000, top=2000, width=1000, height=1000)
        group = FakeShape("grp", st=MSO_SHAPE_TYPE.GROUP,
                          left=0, top=0, width=3000, height=3000,
                          children=[child1, child2])
        comps = _extract_components([group], SLIDE_W, SLIDE_H, _IdCounter())
        # group + 2 children = 3 components
        assert len(comps) == 3
        assert comps[0]["type"] == "group"
        assert {comps[1]["type"], comps[2]["type"]} == {"textbox", "image"}

    def test_empty_container_does_not_crash(self):
        comps = _extract_components([], SLIDE_W, SLIDE_H, _IdCounter())
        assert comps == []


# ---------------------------------------------------------------------------
# (8) Second-template robustness (default Presentation() != bundled template)
# ---------------------------------------------------------------------------
class TestSecondTemplate:
    def test_default_presentation_extracts(self, tmp_path):
        prs = Presentation()  # default blank deck (different master/layouts)
        out = tmp_path / "default.pptx"
        prs.save(str(out))
        schema = extract_schema(str(out))
        r = validate_template_schema(schema)
        assert r.is_valid, r.error_messages()
        assert len(schema["slide_layouts"]) == len(prs.slide_layouts)
        assert isinstance(schema["slide_master"]["components"], list)


# ---------------------------------------------------------------------------
# (9) Negative test — non-PPTX input raises domain error
# ---------------------------------------------------------------------------
class TestNegativePath:
    def test_non_pptx_raises_domain_error(self, tmp_path):
        bad = tmp_path / "not_a_deck.pptx"
        bad.write_bytes(b"this is not a zip file at all")
        with pytest.raises(TemplateExtractionError):
            extract_schema(str(bad))

    def test_missing_file_raises_domain_error(self):
        with pytest.raises(TemplateExtractionError):
            extract_schema("does_not_exist.pptx")
