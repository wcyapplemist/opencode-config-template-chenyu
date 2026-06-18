"""Tests for the resource pipeline: image embedding (#18) + resolvers (#23)."""
import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ppt_builder import generate_ppt_from_data
from resolvers import resolve_slide_data_list


# ============================================================
# Image embedding (#18)
# ============================================================
class TestImageEmbedding:
    def test_slide_with_image_path_embeds_native_picture(self, image_slide_data, output_path):
        generate_ppt_from_data([image_slide_data], output_path=output_path)
        prs = Presentation(output_path)
        pics = [s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(pics) == 1

    def test_picture_is_editable_native_object(self, image_slide_data, output_path):
        generate_ppt_from_data([image_slide_data], output_path=output_path)
        prs = Presentation(output_path)
        pic = next(s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
        # Native pictures have a non-empty image blob (embedded, not linked).
        assert pic.image.blob is not None
        assert len(pic.image.blob) > 0

    def test_missing_image_path_is_graceful(self, output_path):
        data = [{"slide_type": "content_slide", "title": "No Img",
                 "body": "**x** - y", "image_path": "does/not/exist.png", "notes": "n"}]
        generate_ppt_from_data([image for image in data], output_path=output_path)
        prs = Presentation(output_path)
        pics = [s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(pics) == 0  # skipped, not crashed

    @pytest.mark.parametrize("preset", ["full", "half-left", "half-right", "below-title"])
    def test_placement_presets_embed_picture(self, sample_image, preset, output_path):
        data = {"slide_type": "content_slide", "title": "P",
                "body": "**x** - y", "image_path": sample_image,
                "image_position": preset, "notes": "n"}
        generate_ppt_from_data([data], output_path=output_path)
        prs = Presentation(output_path)
        pics = [s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(pics) == 1

    def test_invalid_preset_defaults_gracefully(self, sample_image, output_path):
        data = {"slide_type": "content_slide", "title": "P",
                "body": "**x** - y", "image_path": sample_image,
                "image_position": "nonsense", "notes": "n"}
        generate_ppt_from_data([data], output_path=output_path)
        prs = Presentation(output_path)
        pics = [s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(pics) == 1

    def test_image_size_override(self, sample_image, output_path):
        from pptx.util import Inches
        data = {"slide_type": "content_slide", "title": "P", "body": "**x** - y",
                "image_path": sample_image, "image_position": "full",
                "image_size": {"width": 6, "height": 3}, "notes": "n"}
        generate_ppt_from_data([data], output_path=output_path)
        prs = Presentation(output_path)
        pic = next(s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
        assert pic.width == Inches(6)
        assert pic.height == Inches(3)

    def test_backward_compat_no_image(self, output_path):
        # Slides without image_path render exactly as before.
        data = [{"slide_type": "content_slide", "title": "Plain", "body": "**x** - y"}]
        generate_ppt_from_data(data, output_path=output_path)
        prs = Presentation(output_path)
        pics = [s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(pics) == 0


# ============================================================
# Resolver pipeline (#23) — mocked providers
# ============================================================
class TestImageResolver:
    def test_image_placeholder_resolved_to_path(self, image_placeholder_slide, tmp_path):
        def fetch_fn(query, config):
            return b"\x89PNG\r\n\x1a\nFAKE"

        config = {"image": {"fetch_fn": fetch_fn, "cache_dir": str(tmp_path / "cache")}}
        resolved = resolve_slide_data_list([image_placeholder_slide], config)
        slide = resolved[0]
        assert "image_path" in slide
        assert Path(slide["image_path"]).exists()
        # placeholder keys removed
        assert "image_prompt" not in slide

    def test_image_resolver_graceful_on_failure(self, image_placeholder_slide):
        def fetch_fn(query, config):
            raise RuntimeError("network down")

        config = {"image": {"fetch_fn": fetch_fn}}
        resolved = resolve_slide_data_list([image_placeholder_slide], config)
        # Unchanged (no image_path), build not crashed.
        assert "image_path" not in resolved[0]

    def test_image_resolver_skips_when_unconfigured(self, image_placeholder_slide):
        resolved = resolve_slide_data_list([image_placeholder_slide], config={})
        assert "image_path" not in resolved[0]


class TestIconResolver:
    def test_icon_placeholder_resolved(self, tmp_path):
        lib = tmp_path / "icons"
        lib.mkdir()
        (lib / "growth-up-arrow.svg").write_text("<svg/>")
        slide = {"slide_type": "content_slide", "title": "G", "icon_query": "growth up arrow", "notes": "n"}
        config = {"icon": {"path": str(lib)}}
        resolved = resolve_slide_data_list([slide], config)
        assert "icon_path" in resolved[0]
        assert "icon_query" not in resolved[0]

    def test_icon_no_match_graceful(self, tmp_path):
        lib = tmp_path / "icons"
        lib.mkdir()
        (lib / "camera.svg").write_text("<svg/>")
        slide = {"slide_type": "content_slide", "title": "G", "icon_query": "growth", "notes": "n"}
        resolved = resolve_slide_data_list([slide], {"icon": {"path": str(lib)}})
        assert "icon_path" not in resolved[0]

    def test_icon_no_library_graceful(self):
        slide = {"slide_type": "content_slide", "title": "G", "icon_query": "x", "notes": "n"}
        resolved = resolve_slide_data_list([slide], config={})
        assert "icon_path" not in resolved[0]


class TestChartDataResolver:
    def test_data_query_populates_series(self):
        def search_fn(query, config):
            return {"categories": ["2022", "2023"], "series": [{"name": "R", "values": [10, 12]}],
                    "source": "https://example.com"}
        slide = {"slide_type": "chart_slide", "title": "R", "chart_type": "bar",
                 "data_query": "revenue", "notes": "orig"}
        resolved = resolve_slide_data_list([slide], {"chart_data": {"search_fn": search_fn}})
        out = resolved[0]
        assert out["categories"] == ["2022", "2023"]
        assert out["series"][0]["values"] == [10, 12]
        assert "Data source: https://example.com" in out["notes"]
        assert out["data_source"] == "https://example.com"
        assert "data_query" not in out

    def test_concrete_series_not_overwritten(self):
        def search_fn(query, config):
            return {"categories": ["X"], "series": [{"name": "N", "values": [99]}], "source": "src"}
        slide = {"slide_type": "chart_slide", "title": "R", "chart_type": "bar",
                 "categories": ["A", "B"], "series": [{"name": "S", "values": [1, 2]}],
                 "data_query": "should not overwrite", "notes": "n"}
        resolved = resolve_slide_data_list([slide], {"chart_data": {"search_fn": search_fn}})
        out = resolved[0]
        assert out["categories"] == ["A", "B"]  # unchanged
        assert out["series"][0]["values"] == [1, 2]

    def test_chart_data_graceful_when_unconfigured(self):
        slide = {"slide_type": "chart_slide", "title": "R", "chart_type": "bar",
                 "data_query": "revenue", "notes": "n"}
        resolved = resolve_slide_data_list([slide], config={})
        assert "categories" not in resolved[0]


# ============================================================
# Pipeline-level non-fatal / integration
# ============================================================
class TestResolverPipeline:
    def test_resolver_failure_never_crashes_build(self, image_placeholder_slide):
        # A resolver raising internally is caught; slide passes through unchanged.
        def bad_fetch(query, config):
            raise RuntimeError("boom")

        deck = [image_placeholder_slide,
                {"slide_type": "content_slide", "title": "Other", "notes": "n"}]
        resolved = resolve_slide_data_list(deck, {"image": {"fetch_fn": bad_fetch}})
        assert len(resolved) == 2
        assert resolved[1]["title"] == "Other"

    def test_mixed_deck_passes_through_unresolved(self, mixed_deck):
        # No config: all placeholders unresolved but deck structure intact.
        resolved = resolve_slide_data_list(mixed_deck, config={})
        assert len(resolved) == len(mixed_deck)

    def test_end_to_end_image_render(self, image_placeholder_slide, output_path, tmp_path):
        # Full pipeline: resolve placeholder -> render native picture.
        def fetch_fn(query, config):
            from PIL import Image
            import io
            buf = io.BytesIO()
            Image.new("RGB", (200, 150), (10, 20, 30)).save(buf, format="PNG")
            return buf.getvalue()

        config = {"image": {"fetch_fn": fetch_fn, "cache_dir": str(tmp_path / "cache")}}
        resolved = resolve_slide_data_list([image_placeholder_slide], config)
        generate_ppt_from_data(resolved, output_path=output_path)
        prs = Presentation(output_path)
        pics = [s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(pics) == 1
