"""Verify the prototype output: read back charts and confirm they're native PowerPoint charts."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt

prs = Presentation('../../output/prototype_charts.pptx')
print(f'Slides: {len(prs.slides)}')

for i, slide in enumerate(prs.slides):
    print(f'\n===== Slide {i+1} =====')
    # Find title
    for shape in slide.shapes:
        if shape.has_text_frame and shape == slide.shapes.title:
            print(f'  Title: {slide.shapes.title.text}')
    # Find charts
    for shape in slide.shapes:
        if shape.has_chart:
            chart = shape.chart
            print(f'  Chart: type={chart.chart_type}')
            print(f'    has_title={chart.has_title}, has_legend={chart.has_legend}')
            for j, plot in enumerate(chart.plots):
                print(f'    Plot {j}: series_count={len(plot.series)}, has_data_labels={plot.has_data_labels}')
                for k, series in enumerate(plot.series):
                    print(f'      Series {k}: name="{series.name}"')
                print(f'    Categories: {list(plot.categories)}')
    # Notes
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text
        print(f'  Notes: {notes[:80]}...')
