"""
ppt_builder.py
==============
PPT engine using template.pptx Slide Master layouts with proper placeholders.

Loads the template, adds new slides from named layouts (resolved by name, not
index, so layout reordering does not break it), fills placeholders by type
(TITLE, SUBTITLE, OBJECT), and saves the result.

Layouts are matched by name via ``_LAYOUT_NAME_MAP``; ``template.config.json``
may override the layout name for ``title_slide`` / ``content_slide``.

Usage:
    from ppt_builder import generate_ppt_from_data, DEFAULT_OUTPUT_DIR

    result = generate_ppt_from_data(
        slide_data_list,
        output_path=str(DEFAULT_OUTPUT_DIR / "report.pptx"),
    )
"""

from __future__ import annotations

import json
import logging
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches, Pt

from schema_validator import ValidationError, validate_slide_data_list
from resolvers import resolve_slide_data_list
from template_introspector import get_contract
from contract_adapter import embedded_schema_to_contract
from schema_extractor import read_embedded_schema, TemplateExtractionError

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)

_SCRIPT_DIR = Path(__file__).resolve().parent
TEMPLATES_DIR = _SCRIPT_DIR / "templates"
DEFAULT_OUTPUT_DIR = Path.cwd() / "output"

_TEMPLATE_FILE = TEMPLATES_DIR / "template.pptx"

_TITLE_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
_SUBTITLE_TYPE = PP_PLACEHOLDER.SUBTITLE
_BODY_TYPE = PP_PLACEHOLDER.BODY
_OBJECT_TYPE = PP_PLACEHOLDER.OBJECT
_PICTURE_TYPE = PP_PLACEHOLDER.PICTURE

_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

_LAYOUT_NAME_MAP: Dict[str, List[str]] = {
    "title_slide": ["Title Slide"],
    "closing_slide": ["End"],
    "section_header_slide": ["Section Header"],
    "content_slide": ["Title and Content"],
    "two_content_slide": ["7_Two Content"],
    "comparison_slide": ["Comparison"],
    "content_image_slide": ["Picture with Caption"],
    "chart_slide": ["Blank"],
}

_LAYOUTS_WITH_SUBTITLE = {
    "title_slide", "closing_slide",
}
_LAYOUTS_WITH_BODY = {
    "content_slide", "content_image_slide",
}
_LAYOUTS_WITH_TWO_BODIES = {
    "two_content_slide", "comparison_slide",
}
_LAYOUTS_WITH_CHART = {
    "chart_slide",
}

# Issue #44 (P1): ideal placeholder-composition fingerprint per slide_type.
# Derived from _LAYOUT_NAME_MAP + _LAYOUTS_WITH_* sets. Used by
# _resolve_layout_by_fingerprint() so the engine fills ANY template by
# placeholder composition — layout NAMES become a tie-breaker / fallback, not
# the primary key (DESIGN §6 A2/A3).
_SLIDE_TYPE_FINGERPRINT: Dict[str, List[str]] = {
    "title_slide": ["TITLE", "SUBTITLE"],
    "closing_slide": ["TITLE", "SUBTITLE"],
    "section_header_slide": ["TITLE", "SUBTITLE"],
    "content_slide": ["TITLE", "OBJECT"],
    "two_content_slide": ["TITLE", "OBJECT", "OBJECT"],
    "comparison_slide": ["TITLE", "OBJECT", "OBJECT"],
    "content_image_slide": ["TITLE", "PICTURE"],
    "chart_slide": ["TITLE"],
}

# Slide types whose body/content area is the primary selection concern
# (used for the content_area_in2 tie-break).
_CONTENT_SLIDE_TYPES = {
    "content_slide", "two_content_slide", "comparison_slide", "content_image_slide",
}

# Type-satisfaction relation: which layout placeholder types can SERVE a given
# ideal fingerprint type. A content/body (OBJECT) placeholder is versatile — it
# can host text, pictures, tables or charts — so it satisfies several ideal
# types. This lets e.g. a [TITLE, OBJECT] layout serve a [TITLE, PICTURE] ideal.
_SERVES_LAYOUT: Dict[str, Tuple[str, ...]] = {
    "TITLE": ("TITLE",),
    "SUBTITLE": ("SUBTITLE",),
    "PICTURE": ("PICTURE", "OBJECT"),
    "CHART": ("CHART", "OBJECT"),
    "TABLE": ("TABLE", "OBJECT"),
    "MEDIA": ("MEDIA", "OBJECT"),
    "OBJECT": ("OBJECT",),
}

_CHART_TYPE_MAP: Dict[str, Any] = {
    "bar":                    XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_stacked":            XL_CHART_TYPE.COLUMN_STACKED,
    "bar_horizontal":         XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_horizontal_stacked": XL_CHART_TYPE.BAR_STACKED,
    "pie":                    XL_CHART_TYPE.PIE,
    "pie_exploded":           XL_CHART_TYPE.PIE_EXPLODED,
    "doughnut":               XL_CHART_TYPE.DOUGHNUT,
    "line":                   XL_CHART_TYPE.LINE,
    "line_markers":           XL_CHART_TYPE.LINE_MARKERS,
}

_LEGEND_POSITION_MAP: Dict[str, Any] = {
    "bottom": XL_LEGEND_POSITION.BOTTOM,
    "right":  XL_LEGEND_POSITION.RIGHT,
    "top":    XL_LEGEND_POSITION.TOP,
    "left":   XL_LEGEND_POSITION.LEFT,
}

_CHART_COLORS: List[RGBColor] = [
    RGBColor(0x44, 0x72, 0xC4),
    RGBColor(0xED, 0x7D, 0x31),
    RGBColor(0xFF, 0xC0, 0x00),
    RGBColor(0x5B, 0x9B, 0xD5),
    RGBColor(0x70, 0xAD, 0x47),
    RGBColor(0x95, 0x4F, 0x72),
    RGBColor(0x44, 0x54, 0x6A),
    RGBColor(0xA5, 0xA5, 0xA5),
]

_CHART_FONT_NAME = "Calibri"
_CHART_GRIDLINE_COLOR = RGBColor(0xE7, 0xE6, 0xE6)
_CHART_AXIS_COLOR = RGBColor(0x44, 0x54, 0x6A)
_CHART_TEXT_COLOR = RGBColor(0x44, 0x54, 0x6A)

_CHART_DEFAULT_TYPE = "bar"

_PIE_CHART_TYPES = {"pie", "pie_exploded", "doughnut"}
_BAR_CHART_TYPES = {
    "bar", "bar_stacked", "bar_horizontal", "bar_horizontal_stacked",
}


def _slide_dims_emu(slide: Any) -> Tuple[int, int]:
    """Return ``(width_emu, height_emu)`` of the presentation's slide size."""
    prs = slide.part.package.presentation_part.presentation
    return int(prs.slide_width), int(prs.slide_height)


def _chart_bbox(slide: Any) -> Tuple[int, int, int, int]:
    """Compute a chart bounding box responsive to the slide's actual size.

    The former hard-coded ``_CHART_X/Y/CX/CY`` constants were sized for
    13.33x7.5in widescreen but overflowed on smaller 16:9 templates (e.g.
    10x5.625in).  Margins are proportional so the chart fits any slide size.
    """
    sw, sh = _slide_dims_emu(slide)
    margin_x = max(int(sw * 0.07), int(Inches(0.5)))
    y = max(int(sh * 0.25), int(Inches(1.4)))
    bottom_margin = int(Inches(0.3))
    cx = sw - 2 * margin_x
    cy = sh - y - bottom_margin
    return margin_x, y, cx, cy


# --- Image placement (#18) -------------------------------------------------
_IMAGE_DEFAULT_PRESET = "below-title"
_VALID_IMAGE_PRESETS = {"full", "below-title", "half-left", "half-right"}


def _image_bbox(slide: Any, preset_key: str) -> Dict[str, int]:
    """Compute an image bounding box responsive to the slide's actual size."""
    sw, sh = _slide_dims_emu(slide)
    y = max(int(sh * 0.25), int(Inches(1.4)))
    bottom_margin = int(Inches(0.3))
    cy = sh - y - bottom_margin
    if preset_key in ("half-left", "half-right"):
        margin = int(Inches(0.5))
        half_w = (sw - 3 * margin) // 2
        if preset_key == "half-left":
            return {"x": margin, "y": y, "cx": half_w, "cy": cy}
        return {"x": 2 * margin + half_w, "y": y, "cx": half_w, "cy": cy}
    margin_x = max(int(sw * 0.07), int(Inches(0.5)))
    return {"x": margin_x, "y": y, "cx": sw - 2 * margin_x, "cy": cy}


def _normalize_layout_name(name: str) -> str:
    return re.sub(r"^\d+_", "", name).strip().lower()


def _build_layout_index(prs: Presentation) -> Tuple[Dict[str, Any], Dict[str, Any]]:
    exact: Dict[str, Any] = {}
    normalized: Dict[str, Any] = {}
    for layout in prs.slide_layouts:
        nm = layout.name
        exact[nm.lower()] = layout
        norm_key = _normalize_layout_name(nm)
        normalized.setdefault(norm_key, layout)
    return exact, normalized


def _resolve_layout(
    candidate_names: List[str],
    exact: Dict[str, Any],
    normalized: Dict[str, Any],
) -> Optional[Any]:
    for cand in candidate_names:
        if cand.lower() in exact:
            return exact[cand.lower()]
    for cand in candidate_names:
        key = _normalize_layout_name(cand)
        if key in normalized:
            return normalized[key]
    return None


# ---------------------------------------------------------------------------
# Issue #44 (P1): fingerprint-based layout resolution
# ---------------------------------------------------------------------------
def _composition_diff(ideal: List[str], layout_fp: List[str]) -> Tuple[int, int]:
    """Return ``(missing, extra)`` between an ideal fingerprint and a layout's.

    ``missing`` = ideal types that no layout placeholder can serve (each ideal
    type consumes one distinct *serving* placeholder). ``extra`` = layout
    placeholders left unconsumed after satisfying the ideal. A content/body
    (OBJECT) placeholder is versatile — it can also serve PICTURE/TABLE/CHART.
    """
    avail: Dict[str, int] = {}
    for t in layout_fp:
        avail[t] = avail.get(t, 0) + 1
    matched = 0
    # Satisfy non-OBJECT ideal types first so OBJECT placeholders are reserved.
    for it in ideal:
        if it == "OBJECT":
            continue
        for lt in _SERVES_LAYOUT.get(it, (it,)):
            if avail.get(lt, 0) > 0:
                avail[lt] -= 1
                matched += 1
                break
    object_need = sum(1 for it in ideal if it == "OBJECT")
    served = min(object_need, avail.get("OBJECT", 0))
    matched += served
    missing = len(ideal) - matched
    extra = len(layout_fp) - matched
    return missing, extra


def _name_affinity(layout_name: str, candidate_names: List[str]) -> int:
    """``2`` = exact (case-insensitive) name match, ``1`` = normalized, ``0`` = none."""
    cname = layout_name.lower()
    if any(cname == c.lower() for c in candidate_names):
        return 2
    norm = _normalize_layout_name(layout_name)
    if any(norm == _normalize_layout_name(c) for c in candidate_names):
        return 1
    return 0


_LAYOUT_TYPES_NEEDING_SIDEBYSIDE = {"two_content_slide", "comparison_slide"}


def _content_placeholders_stacked(layout_contract: Dict[str, Any]) -> int:
    """Return 1 if this layout's content placeholders are vertically stacked,
    0 if horizontally separated (side-by-side).

    Used as a geometric tie-breaker so two_content/comparison slides prefer
    side-by-side layouts over vertically stacked ones.
    """
    phs = [p for p in layout_contract.get("placeholders", [])
           if p.get("type") == "OBJECT"]
    if len(phs) < 2:
        return 0
    lefts = sorted(p.get("left_in", 0) for p in phs[:2])
    return 0 if (lefts[1] - lefts[0]) > 1.0 else 1


def _resolve_layout_by_fingerprint(
    slide_type: str,
    contract: Dict[str, Any],
) -> Tuple[Optional[int], Optional[str]]:
    """Match ``slide_type`` to the best contract layout by composition.

    Returns ``(layout_index, None)`` on success, or ``(None, reason)`` on
    degradation (no layout can satisfy the ideal composition). Among
    composition-compatible layouts, ranking is: name affinity (highest) → fewest
    surplus placeholders → largest content area → lowest index. Names are a
    tie-breaker, not the primary key (DESIGN §6 A2, issue #44).
    """
    ideal = _SLIDE_TYPE_FINGERPRINT.get(slide_type)
    layouts = (contract or {}).get("layouts", [])
    if ideal is None:
        return None, f"no fingerprint defined for slide_type '{slide_type}'"
    if not layouts:
        return None, "contract has no layouts"

    candidate_names = _LAYOUT_NAME_MAP.get(slide_type, [])
    need_side_by_side = slide_type in _LAYOUT_TYPES_NEEDING_SIDEBYSIDE
    scored: List[Tuple[int, int, int, int, float, int, str]] = []
    for L in layouts:
        missing, extra = _composition_diff(ideal, L.get("fingerprint", []))
        affinity = _name_affinity(L.get("name", ""), candidate_names)
        stacked = _content_placeholders_stacked(L) if need_side_by_side else 0
        # sort key: (missing, -affinity, extra, stacked, -area, index)
        # — min() picks the lowest missing, then highest affinity, fewest
        # extras, side-by-side over stacked, largest area, lowest index.
        scored.append((
            missing, -affinity, extra, stacked,
            -float(L.get("content_area_in2", 0)), L["index"], L.get("name", ""),
        ))

    full = [s for s in scored if s[0] == 0]
    if not full:
        best = min(scored)
        return None, (
            f"no layout satisfies fingerprint {ideal} "
            f"(closest: '{best[6]}' missing {best[0]})"
        )
    best = min(full)
    return best[5], None


def _select_layout(
    slide_type: str,
    contract: Optional[Dict[str, Any]],
    config: Dict[str, Any],
    prs: Presentation,
    exact_idx: Dict[str, Any],
    norm_idx: Dict[str, Any],
    page_num: int,
) -> Optional[Any]:
    """Resolve a slide_type to a concrete ``SlideLayout`` (issue #44).

    Precedence: config pin (``<slide_type>_layout``) → fingerprint match →
    name-based fallback → degradation (skip + warn). Without a contract the path
    is the original name-based matching, so behaviour is backward compatible.
    """
    if slide_type not in _SLIDE_TYPE_FINGERPRINT and slide_type not in _LAYOUT_NAME_MAP:
        logger.warning("Page %d: unknown slide_type '%s', skipped", page_num, slide_type)
        return None

    # 1. Config pin — explicit layout name (highest precedence; all 8 types).
    pinned = config.get(f"{slide_type}_layout")
    if pinned:
        layout = _resolve_layout([pinned], exact_idx, norm_idx)
        if layout is not None:
            return layout
        logger.warning(
            "Page %d: config pin '%s' not found; falling back", page_num, pinned)

    # 2. Fingerprint match (contract-aware; names are a tie-breaker).
    if contract:
        idx, reason = _resolve_layout_by_fingerprint(slide_type, contract)
        if idx is not None:
            return prs.slide_layouts[idx]
        logger.warning(
            "Page %d: fingerprint degradation for '%s': %s",
            page_num, slide_type, reason)

    # 3. Name-based fallback (backward-compatible safety net).
    candidates = _LAYOUT_NAME_MAP.get(slide_type)
    if candidates:
        layout = _resolve_layout(candidates, exact_idx, norm_idx)
        if layout is not None:
            return layout

    # 4. Degradation: nothing usable.
    logger.warning(
        "Page %d: no layout matched for slide_type '%s', skipped", page_num, slide_type)
    return None


def servable_slide_types(contract: Dict[str, Any]) -> Dict[str, Dict[str, Any]]:
    """Report which engine slide_types a template's contract can serve (#45).

    Used by the content/outline stage to constrain itself to layouts the
    template actually provides (never emit a ``slide_type`` that would degrade).
    For each of the 8 slide types, returns ``{"available": bool, ...}`` — when
    available, the selected layout name/index; when not, the degradation reason.
    """
    layouts = (contract or {}).get("layouts", [])
    report: Dict[str, Dict[str, Any]] = {}
    for slide_type in _SLIDE_TYPE_FINGERPRINT:
        idx, reason = _resolve_layout_by_fingerprint(slide_type, contract)
        if idx is not None and 0 <= idx < len(layouts):
            report[slide_type] = {
                "available": True,
                "layout": layouts[idx].get("name", ""),
                "index": idx,
                "content_area_in2": layouts[idx].get("content_area_in2", 0),
            }
        else:
            report[slide_type] = {"available": False, "reason": reason}
    return report


def _load_config() -> Dict[str, Any]:
    config_path = TEMPLATES_DIR / "template.config.json"
    if config_path.exists():
        with open(config_path, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def _remove_all_slides(prs: Presentation) -> int:
    count = len(prs.slides)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].attrib.get(f"{{{_REL_NS}}}id")
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    return count


def _find_placeholder(slide: Any, ph_type: Any) -> Optional[Any]:
    for ph in slide.placeholders:
        if ph.placeholder_format.type == ph_type:
            return ph
    return None


def _find_placeholders(slide: Any, ph_type: Any) -> List[Any]:
    return [ph for ph in slide.placeholders if ph.placeholder_format.type == ph_type]


def _find_title_placeholder(slide: Any) -> Optional[Any]:
    for ph_type in _TITLE_TYPES:
        ph = _find_placeholder(slide, ph_type)
        if ph:
            return ph
    return None


def _find_body_placeholder(slide: Any) -> Optional[Any]:
    ph = _find_placeholder(slide, _BODY_TYPE)
    if ph:
        return ph
    objects = _find_placeholders(slide, _OBJECT_TYPE)
    return objects[0] if objects else None


def _set_text(shape: Any, text: str) -> bool:
    if not shape or not shape.has_text_frame:
        return False
    try:
        tf = shape.text_frame
        tf.clear()
        tf.paragraphs[0].text = text
        return True
    except Exception as exc:
        logger.warning("Failed to set text: %s", exc)
        return False


def _parse_line(line: str) -> Tuple[str, str]:
    clean = re.sub(r"\*\*", "", line.strip())
    if not clean:
        return ("", "")
    for sep in [" \u2014 ", " - ", ": "]:
        if sep in clean:
            parts = clean.split(sep, 1)
            return (parts[0].strip(), parts[1].strip())
    return (clean, "")


def _set_notes(slide: Any, notes_text: str) -> bool:
    text = (notes_text or "").strip()
    if not text:
        return False
    try:
        slide.notes_slide.notes_text_frame.text = text
        return True
    except Exception as exc:
        logger.warning("Failed to set notes: %s", exc)
        return False


def _set_body_text(shape: Any, text: str) -> bool:
    if not shape or not shape.has_text_frame:
        return False
    try:
        lines = [l for l in text.split("\n") if l.strip()]
        tf = shape.text_frame
        tf.clear()

        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            title_part, desc_part = _parse_line(line)

            if title_part:
                run = p.add_run()
                run.text = title_part
                run.font.bold = True
                run.font.size = Pt(14)
            if desc_part:
                run = p.add_run()
                run.text = f" \u2014 {desc_part}" if title_part else desc_part
                run.font.size = Pt(12)

        return True
    except Exception as exc:
        logger.warning("Failed to set body text: %s", exc)
        return False


def _apply_series_colors(chart: Any, chart_type_key: str) -> None:
    is_pie = chart_type_key in _PIE_CHART_TYPES
    try:
        plot = chart.plots[0]
        if is_pie:
            for idx, point in enumerate(plot.series[0].points):
                color = _CHART_COLORS[idx % len(_CHART_COLORS)]
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = color
        else:
            for idx, series in enumerate(plot.series):
                color = _CHART_COLORS[idx % len(_CHART_COLORS)]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = color
                if chart_type_key.startswith("line"):
                    series.format.line.color.rgb = color
                    series.format.line.width = Pt(2.5)
    except Exception as exc:
        logger.warning("Failed to apply series colors: %s", exc)


def _add_chart_to_slide(slide: Any, slide_data: Dict[str, Any]) -> bool:
    chart_type_key = slide_data.get("chart_type", _CHART_DEFAULT_TYPE)
    if chart_type_key not in _CHART_TYPE_MAP:
        logger.warning(
            "Unknown chart_type '%s', defaulting to '%s'",
            chart_type_key, _CHART_DEFAULT_TYPE,
        )
        chart_type_key = _CHART_DEFAULT_TYPE

    categories = slide_data.get("categories", [])
    series_list = slide_data.get("series", [])

    if not categories or not series_list:
        logger.warning(
            "Chart slide missing categories or series, skipping chart"
        )
        return False

    chart_data = CategoryChartData()
    chart_data.categories = list(categories)
    for s in series_list:
        name = s.get("name", "")
        values = list(s.get("values", []))
        chart_data.add_series(name, values)

    xl_type = _CHART_TYPE_MAP[chart_type_key]
    try:
        cx_chart_x, cx_chart_y, cx_chart_cx, cx_chart_cy = _chart_bbox(slide)
        graphic_frame = slide.shapes.add_chart(
            xl_type,
            cx_chart_x, cx_chart_y, cx_chart_cx, cx_chart_cy,
            chart_data,
        )
    except Exception as exc:
        logger.error("Failed to create chart: %s", exc)
        return False

    chart = graphic_frame.chart
    options = slide_data.get("chart_options", {})

    chart.has_title = False
    chart.font.name = _CHART_FONT_NAME
    chart.font.size = Pt(11)

    legend_pos_key = options.get("legend_position", "bottom")
    if legend_pos_key == "none":
        chart.has_legend = False
    else:
        chart.has_legend = True
        chart.legend.position = _LEGEND_POSITION_MAP.get(
            legend_pos_key, XL_LEGEND_POSITION.BOTTOM,
        )
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(11)
        chart.legend.font.name = _CHART_FONT_NAME
        chart.legend.font.color.rgb = _CHART_TEXT_COLOR

    is_pie = chart_type_key in _PIE_CHART_TYPES
    is_bar = chart_type_key in _BAR_CHART_TYPES

    try:
        plot = chart.plots[0]
        show_labels = options.get("show_data_labels", True)
        plot.has_data_labels = show_labels
        if show_labels:
            labels = plot.data_labels
            labels.font.size = Pt(10)
            labels.font.name = _CHART_FONT_NAME
            labels.font.color.rgb = _CHART_TEXT_COLOR
            if is_pie:
                labels.show_percentage = True
                labels.show_value = False
                labels.show_category_name = False
                labels.number_format = "0%"
            else:
                labels.show_value = True
                labels.show_percentage = False
                labels.number_format = options.get("value_format", "#,##0.0")
                if is_bar:
                    labels.position = XL_LABEL_POSITION.OUTSIDE_END
    except Exception as exc:
        logger.warning("Failed to set data labels: %s", exc)

    if not is_pie:
        try:
            val_axis = chart.value_axis
            val_axis.has_major_gridlines = True
            val_axis.major_gridlines.format.line.color.rgb = _CHART_GRIDLINE_COLOR
            val_axis.major_gridlines.format.line.width = Pt(0.75)
            val_axis.tick_labels.font.size = Pt(10)
            val_axis.tick_labels.font.name = _CHART_FONT_NAME
            val_axis.tick_labels.font.color.rgb = _CHART_AXIS_COLOR
            val_axis.format.line.color.rgb = _CHART_AXIS_COLOR
            val_axis.tick_labels.number_format = options.get("y_axis_format", "#,##0.0")
            if options.get("y_axis_min") is not None:
                val_axis.minimum_scale = options["y_axis_min"]
            if options.get("y_axis_max") is not None:
                val_axis.maximum_scale = options["y_axis_max"]
            if options.get("y_axis_major_unit") is not None:
                val_axis.major_unit = options["y_axis_major_unit"]
            if options.get("y_axis_title"):
                val_axis.has_title = True
                val_axis.axis_title.text_frame.text = options["y_axis_title"]
                val_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
                val_axis.axis_title.text_frame.paragraphs[0].font.name = _CHART_FONT_NAME
                val_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = _CHART_AXIS_COLOR

            cat_axis = chart.category_axis
            cat_axis.tick_labels.font.size = Pt(10)
            cat_axis.tick_labels.font.name = _CHART_FONT_NAME
            cat_axis.tick_labels.font.color.rgb = _CHART_AXIS_COLOR
            cat_axis.format.line.color.rgb = _CHART_AXIS_COLOR
            if options.get("x_axis_title"):
                cat_axis.has_title = True
                cat_axis.axis_title.text_frame.text = options["x_axis_title"]
                cat_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
                cat_axis.axis_title.text_frame.paragraphs[0].font.name = _CHART_FONT_NAME
                cat_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = _CHART_AXIS_COLOR
        except Exception as exc:
            logger.warning("Failed to set axis options: %s", exc)

    _apply_series_colors(chart, chart_type_key)

    logger.info(
        "  Chart: type=%s, categories=%d, series=%d",
        chart_type_key, len(categories), len(series_list),
    )
    return True


def _find_picture_placeholder(slide: Any) -> Optional[Any]:
    for ph in slide.placeholders:
        if ph.placeholder_format.type == _PICTURE_TYPE:
            return ph
    return None


def _add_image_to_slide(slide: Any, slide_data: Dict[str, Any]) -> bool:
    """Embed a native, editable PowerPoint picture from ``image_path`` (#18).

    Placement order:
      1. If the layout has a PICTURE placeholder, fill it natively.
      2. Otherwise place using a named preset (``image_position``) or an
         explicit ``image_size`` override, in the free space below the title.

    Images are embedded (not linked) so the PPTX stays self-contained and the
    picture remains fully editable in PowerPoint.
    """
    image_path = slide_data.get("image_path")
    if not image_path:
        return False

    p = Path(image_path)
    if not p.exists():
        logger.warning("image_path not found: %s, skipping image", image_path)
        return False

    preset_key = slide_data.get("image_position", _IMAGE_DEFAULT_PRESET)
    if preset_key not in _VALID_IMAGE_PRESETS:
        logger.warning(
            "Unknown image_position '%s', defaulting to '%s'",
            preset_key, _IMAGE_DEFAULT_PRESET,
        )
        preset_key = _IMAGE_DEFAULT_PRESET

    try:
        pic_ph = _find_picture_placeholder(slide)
        if pic_ph is not None:
            try:
                pic_ph.insert_picture(str(p))
                logger.info("  Image (placeholder): %s", p.name)
                return True
            except Exception:
                # Fall back to free placement at the placeholder's frame box.
                box = {
                    "x": pic_ph.left, "y": pic_ph.top,
                    "cx": pic_ph.width, "cy": pic_ph.height,
                }
                slide.shapes.add_picture(str(p), box["x"], box["y"], box["cx"], box["cy"])
                logger.info("  Image (placeholder frame): %s", p.name)
                return True

        box = _image_bbox(slide, preset_key)
        cx, cy = box["cx"], box["cy"]
        size = slide_data.get("image_size")
        if isinstance(size, dict):
            if size.get("width"):
                cx = Inches(size["width"])
            if size.get("height"):
                cy = Inches(size["height"])
        slide.shapes.add_picture(str(p), box["x"], box["y"], cx, cy)
        logger.info("  Image (%s): %s", preset_key, p.name)
        return True
    except Exception as exc:
        logger.error("Failed to embed image '%s': %s", image_path, exc)
        return False


_DEFAULT_CLOSING_NOTES = (
    'KEY MESSAGE: Thank you — close warmly and open the floor for questions.\n'
    '"Thank you all for your time today."\n'
    'Pause. Make eye contact across the room.\n'
    '"I hope this gave you a clear picture of where we are and where we are headed."\n'
    'TRANSITION: "I would love to take any questions you have."\n'
    'COACHING: Warm, unhurried close. Be ready for: "Can you share the deck?" '
    '— yes, I will send it after.'
)


def _ensure_default_closing(
    slide_data_list: List[Dict[str, Any]],
) -> List[Dict[str, Any]]:
    # #40 / US-6: append a default Thank-You closing when the deck is large
    # enough and does not already end on a closing_slide. Non-destructive —
    # returns a new list, leaving the caller's list untouched.
    if not isinstance(slide_data_list, list) or len(slide_data_list) < 3:
        return slide_data_list
    last_type = (slide_data_list[-1] or {}).get("slide_type", "")
    if last_type == "closing_slide":
        return slide_data_list
    closing = {
        "slide_type": "closing_slide",
        "title": "Thank You",
        "notes": _DEFAULT_CLOSING_NOTES,
    }
    logger.info("Auto-appending default closing slide (default_closing=True)")
    return list(slide_data_list) + [closing]


def get_render_contract(template_path: str) -> Dict[str, Any]:
    """Return the render contract for ``template_path`` (US-4.1).

    Prefers the embedded ``ppt/template_schema.json`` (via the
    :mod:`contract_adapter` bridge) and falls back to the mtime-cached sidecar
    introspection contract (:func:`get_contract`). Provenance is tagged on the
    returned dict as ``_source ∈ {"embedded", "sidecar"}`` and logged.

    Failure handling (architecture review M4):
      - embedded JSON absent (legacy template) -> silent sidecar fallback.
      - embedded JSON malformed -> ``read_embedded_schema`` already warns;
        sidecar fallback.
      - corrupt/non-zip input -> ``TemplateExtractionError`` caught here, warned,
        sidecar fallback.

    May raise if the sidecar contract itself fails (the caller's try/except then
    degrades to name-based layout matching — backward compatible).
    """
    try:
        schema = read_embedded_schema(template_path)
    except TemplateExtractionError as exc:
        logger.warning(
            "Embedded schema unreadable for %s (%s); falling back to sidecar contract",
            template_path, exc,
        )
        schema = None
    except Exception as exc:  # defensive — never block the render
        logger.warning(
            "Embedded schema read failed for %s (%s); sidecar fallback",
            template_path, exc,
        )
        schema = None

    if schema is not None:
        try:
            contract = embedded_schema_to_contract(schema)
            contract["_source"] = "embedded"
            logger.info(
                "Render contract (embedded): %d layouts, ratio %s",
                len(contract.get("layouts", [])),
                contract.get("slide_size", {}).get("ratio", "?"),
            )
            return contract
        except Exception as exc:
            logger.warning(
                "Embedded schema -> contract failed for %s (%s); sidecar fallback",
                template_path, exc,
            )

    contract = get_contract(str(template_path))
    contract["_source"] = "sidecar"
    logger.info(
        "Render contract (sidecar): %d layouts, ratio %s",
        len(contract.get("layouts", [])),
        contract.get("slide_size", {}).get("ratio", "?"),
    )
    return contract


def generate_ppt_from_data(
    slide_data_list: List[Dict[str, Any]],
    template_path: Optional[str] = None,
    output_path: str = "output.pptx",
    prompt_text: str = "",
    validate: bool = True,
    strict: bool = False,
    cleanup_temp: bool = True,
    resolve_placeholders: bool = True,
    default_closing: bool = True,
    config_overrides: Optional[Dict[str, str]] = None,
) -> str:
    # #37: resolve resource placeholders (data_query) into concrete assets
    # BEFORE validation, so the validator sees materialized data.
    # Graceful no-op when resolver.config.json is absent.
    if resolve_placeholders:
        slide_data_list = resolve_slide_data_list(slide_data_list)

    # #40 / US-6: guarantee a Thank-You closing slide on decks of N >= 3.
    if default_closing:
        slide_data_list = _ensure_default_closing(slide_data_list)

    # Phase 1 Track A: defensive validation. Catches malformed input with a
    # clear ValidationError instead of a cryptic crash in the render loop.
    if validate:
        result = validate_slide_data_list(slide_data_list, strict=strict)
        for msg in result.warning_messages():
            logger.warning("Validation: %s", msg)
        # Strict mode (agent pre-flight gate): any schema error is fatal.
        if strict and not result.is_valid:
            raise ValidationError(result.errors)
        # Non-strict mode: surface per-slide schema errors as warnings for
        # visibility, but keep the engine's existing graceful-degradation
        # behaviour (skip slide / default chart / skip chart). Only abort on
        # unrecoverable top-level structure, which would otherwise crash the
        # render loop cryptically.
        if not strict:
            for msg in result.error_messages():
                logger.warning("Validation (degraded): %s", msg)
        if not isinstance(slide_data_list, list):
            raise ValidationError(
                result.errors if result.errors
                else "slide_data_list must be a JSON array"
            )

    template = Path(template_path) if template_path and template_path != "auto" else _TEMPLATE_FILE
    output = Path(output_path)

    if not template.exists():
        raise FileNotFoundError(f"Template not found: {template}")

    # #46 (P3): state machine ① — discard any derived template_new.pptx left
    # from a prior run so the base template.pptx is re-evaluated fresh each
    # request. Inline (no cross-skill import); the full lifecycle lives in the
    # template-modifier-skill and is wired by P4 when cloning is implemented.
    _derived = template.with_name(template.stem + "_new" + template.suffix)
    if _derived.exists():
        try:
            _derived.unlink()
            logger.info("Discarded leftover derived template: %s", _derived.name)
        except OSError as exc:  # pragma: no cover - defensive
            logger.debug("Could not delete leftover %s: %s", _derived, exc)

    if not output.is_absolute():
        output = DEFAULT_OUTPUT_DIR / output
    output.parent.mkdir(parents=True, exist_ok=True)

    config = _load_config()
    # #47 (P4): merge caller-supplied layout overrides (e.g. cloned-layout pins
    # from template-modifier-skill's resolve_and_clone). Caller overrides win.
    if config_overrides:
        config = {**config, **config_overrides}

    logger.info("Loading template: %s", template.name)
    prs = Presentation(str(template))
    logger.info("Template: %d slides, %d layouts", len(prs.slides), len(prs.slide_layouts))

    # #43 (P0): auto-introspect the template into a JSON contract before render.
    # US-4.1: prefer the embedded JSON (via the adapter); fall back to the
    # mtime-cached sidecar contract. Non-fatal: on any failure the engine falls
    # back to name-based layout matching (backward compatible). Provenance is
    # tagged on the contract as ``_source`` and logged inside get_render_contract.
    contract = None
    try:
        contract = get_render_contract(str(template))
    except Exception as exc:  # pragma: no cover - defensive; never block render
        logger.warning("Template contract unavailable (%s); using name matching", exc)

    removed = _remove_all_slides(prs)
    logger.info("Cleared %d example slides", removed)

    exact_idx, norm_idx = _build_layout_index(prs)

    for page_num, slide_data in enumerate(slide_data_list, start=1):
        slide_type = slide_data.get("slide_type", "")

        # Resolve layout: config pin → fingerprint match → name fallback (#44).
        layout = _select_layout(
            slide_type, contract, config, prs, exact_idx, norm_idx, page_num
        )
        if layout is None:
            continue  # degradation warning already logged

        layout_idx = prs.slide_layouts.index(layout)
        try:
            slide = prs.slides.add_slide(layout)
            logger.info("Page %d: added slide from layout[%d] '%s'", page_num, layout_idx, layout.name)

            title_text = slide_data.get("title", "")

            # Always try to fill title placeholder
            if title_text:
                title_ph = _find_title_placeholder(slide)
                if title_ph:
                    _set_text(title_ph, title_text)
                    logger.info("  Title: \"%s\"", title_text)

            # Fill subtitle (for title, agenda, closing, section-header-sub)
            if slide_type in _LAYOUTS_WITH_SUBTITLE:
                subtitle_text = slide_data.get("subtitle", "")
                if subtitle_text:
                    sub_ph = _find_placeholder(slide, _SUBTITLE_TYPE)
                    if sub_ph:
                        _set_text(sub_ph, subtitle_text)
                        logger.info("  Subtitle: \"%s\"", subtitle_text[:50])

            # Fill body text (for content slides)
            if slide_type in _LAYOUTS_WITH_BODY:
                body_text = slide_data.get("body", "")
                if body_text:
                    body_ph = _find_body_placeholder(slide)
                    if body_ph:
                        _set_body_text(body_ph, body_text)
                        logger.info("  Body: %d lines", len([l for l in body_text.split("\n") if l.strip()]))

            # Fill two body areas (for two-content slides)
            if slide_type in _LAYOUTS_WITH_TWO_BODIES:
                body_left = slide_data.get("body_left", "")
                body_right = slide_data.get("body_right", "")
                objects = _find_placeholders(slide, _OBJECT_TYPE)
                # BODY placeholders serve the same role as OBJECT (the
                # introspector normalizes BODY→OBJECT in the contract, but
                # _find_placeholders uses the raw python-pptx type). Fall
                # back so BODY-based templates fill correctly.
                if len(objects) < 2:
                    body_phs = _find_placeholders(slide, _BODY_TYPE)
                    if len(body_phs) > len(objects):
                        objects = body_phs
                if len(objects) >= 2:
                    if body_left:
                        _set_body_text(objects[0], body_left)
                        logger.info("  Body-left: %d lines", len([l for l in body_left.split("\n") if l.strip()]))
                    if body_right:
                        _set_body_text(objects[1], body_right)
                        logger.info("  Body-right: %d lines", len([l for l in body_right.split("\n") if l.strip()]))
                elif len(objects) == 1 and (body_left or body_right):
                    _set_body_text(objects[0], body_left or body_right)
                    logger.warning(
                        "  Two-content slide has only 1 content placeholder; "
                        "body_left/body_right merged into one")
                elif body_left or body_right:
                    logger.warning(
                        "  Two-content slide has no content placeholders; "
                        "body_left/body_right dropped")

            # Add chart (for chart slides)
            if slide_type in _LAYOUTS_WITH_CHART:
                _add_chart_to_slide(slide, slide_data)

            # Embed image (any slide carrying image_path) — #18
            if slide_data.get("image_path"):
                _add_image_to_slide(slide, slide_data)

            # Fill speaker notes (must be English; only visible in Presenter View)
            notes_text = slide_data.get("notes", "")
            if _set_notes(slide, notes_text):
                logger.info("  Notes: %d chars", len(notes_text))

        except Exception as exc:
            logger.error("Page %d failed: %s", page_num, exc)

    prs.save(str(output))
    logger.info("Saved: %s (%d slides)", output.resolve(), len(prs.slides))

    # Auto-cleanup pipeline temp artifacts (outline checkpoints, agent-written
    # temp JSON) so they never accumulate on disk. Lazy import + try/except keeps
    # cleanup from ever affecting a successful render. Pass cleanup_temp=False to
    # retain them (e.g. while debugging a failed run).
    if cleanup_temp:
        try:
            from outline_store import cleanup_all
            removed = cleanup_all()
            if removed:
                logger.info("Cleaned up %d temp artifact(s)", removed)
        except Exception as exc:  # cleanup must never break a successful render
            logger.debug("Temp cleanup skipped: %s", exc)
    return str(output.resolve())


def main() -> None:
    mock: List[Dict[str, Any]] = [
        {
            "slide_type": "title_slide",
            "title": "AI Empowering Finance",
            "subtitle": "2026 Q1",
            "notes": (
                "KEY MESSAGE: Open with energy — set the stakes in one line.\n"
                "\"Hold the slide for two seconds before you speak.\"\n"
                "\"Good [morning/afternoon], I'm [Name]. Today I want to show you how AI is already transforming finance — not in theory, but in the numbers.\"\n"
                "Pause. Let the tagline land.\n"
                "\"We'll walk through where it delivers the clearest ROI today.\"\n"
                "TRANSITION: \"Let me start with the core scenarios.\"\n"
                "COACHING: Eye contact, confident. Do not read the slide. Be ready for: \"Is this hype or real?\" — lead with the 80 percent figure."
            ),
        },
        {
            "slide_type": "content_slide",
            "title": "Core AI Scenarios",
            "body": (
                "**Automated Reporting** \u2014 RPA tools auto-generate monthly reports, cutting manual effort by 80%\n"
                "**Smart Reconciliation** \u2014 AI matches bank transactions at 99.5% accuracy\n"
                "**Fraud Detection** \u2014 Real-time anomaly detection with automated alerts\n"
                "**Tax Optimization** \u2014 ML identifies savings opportunities across tax structures"
            ),
            "notes": (
                "KEY MESSAGE: Four high-impact scenarios where AI already delivers measurable ROI.\n"
                "\"Let's make this concrete. These aren't edge cases — this is everyday finance.\"\n"
                "\"Automated reporting alone removes eighty percent of the manual effort behind every monthly close.\"\n"
                "Pause. Let the number land.\n"
                "\"Smart reconciliation now matches transactions at ninety-nine-point-five percent accuracy, and fraud detection flags anomalies in real time.\"\n"
                "\"Ask your CFO: how much would one missed discrepancy cost?\"\n"
                "TRANSITION: \"Here is how we roll this out.\"\n"
                "COACHING: Matter-of-fact tone, don't over-sell. Be ready for: \"What about false positives?\" — answer: tuned thresholds, human-in-the-loop review."
            ),
        },
        {
            "slide_type": "content_slide",
            "title": "Roadmap",
            "body": (
                "**Phase 1: Pilot** \u2014 Deploy in 2 business units by Q2\n"
                "**Phase 2: Scale** \u2014 Expand to all departments by Q4\n"
                "**Phase 3: Full Deployment** \u2014 Organization-wide adoption by 2027"
            ),
            "notes": (
                "KEY MESSAGE: A phased, low-risk rollout — pilot, scale, then full adoption.\n"
                "\"We don't boil the ocean. We pilot in two units first, prove the numbers, then scale.\"\n"
                "\"By Q4 every department is on board, and full organisation-wide adoption lands in 2027.\"\n"
                "Walk the three phases left to right.\n"
                "TRANSITION: Open for questions.\n"
                "COACHING: Keep it tight, end with confidence. Be ready for: \"What could delay Phase 2?\" — answer: only change-management, never the technology."
            ),
        },
    ]

    print("Test: template.pptx (placeholder-based)")
    result = generate_ppt_from_data(
        mock,
        output_path=str(DEFAULT_OUTPUT_DIR / "test_template.pptx"),
    )
    print(f"Output: {result}")


if __name__ == "__main__":
    main()
