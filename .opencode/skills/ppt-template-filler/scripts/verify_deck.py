"""Verify the generated construction deck PPT."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

prs = Presentation('output/202606150807.pptx')
print(f'Slides: {len(prs.slides)}')

slide_types_used = set()
chart_count = 0
notes_count = 0

for i, slide in enumerate(prs.slides):
    shapes_count = len(slide.shapes)
    layout_name = slide.slide_layout.name

    title = ""
    for shape in slide.shapes:
        if shape.has_text_frame and shape == slide.shapes.title:
            title = slide.shapes.title.text
            break

    has_chart = any(s.has_chart for s in slide.shapes)
    has_notes = slide.has_notes_slide and bool(slide.notes_slide.notes_text_frame.text.strip())

    if has_chart:
        chart_count += 1
        for s in slide.shapes:
            if s.has_chart:
                ct = s.chart.chart_type
                print(f'  Slide {i+1}: [{layout_name}] "{title}" CHART type={ct}')
    else:
        print(f'  Slide {i+1}: [{layout_name}] "{title}" ({shapes_count} shapes)')

    slide_types_used.add(layout_name)
    if has_notes:
        notes_count += 1

print(f'\nLayouts used ({len(slide_types_used)}):')
for lt in sorted(slide_types_used):
    print(f'  - {lt}')

print(f'\nCharts: {chart_count}')
print(f'Slides with notes: {notes_count}/{len(prs.slides)}')
print(f'\nAll checks: {"PASS" if len(prs.slides) == 16 and chart_count == 2 and notes_count == 16 and len(slide_types_used) >= 7 else "CHECK"}')
