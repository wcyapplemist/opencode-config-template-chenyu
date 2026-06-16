"""
Prototype: validate that python-pptx can generate native charts on a Blank layout.
Creates bar, pie, and line charts on template slides and saves to output/.
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
sys.path.insert(0, '.')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

TEMPLATE = 'templates/template.pptx'
OUTPUT = '../../output/prototype_charts.pptx'

prs = Presentation(TEMPLATE)

# Remove all existing slides (same as engine does)
from pptx.oxml.ns import qn
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(qn('r:id'))
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# Use layout [9] "Blank" - has TITLE placeholder
blank_layout = prs.slide_layouts[9]

# ============================================================
# Chart 1: BAR CHART (COLUMN_CLUSTERED)
# ============================================================
slide1 = prs.slides.add_slide(blank_layout)

# Fill title
for ph in slide1.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text_frame.text = "Global Construction Tech Market (USD Billion)"
        break

chart_data = CategoryChartData()
chart_data.categories = ['2020', '2021', '2022', '2023', '2024', '2025', '2026']
chart_data.add_series('Market Size', (8.5, 11.2, 14.8, 19.5, 25.1, 31.7, 39.4))

# Position: left=0.92in, top=2.0in, w=11.5in, h=4.5in
graphic_frame = slide1.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.92), Inches(2.0), Inches(11.5), Inches(4.5),
    chart_data
)
chart = graphic_frame.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

# Style the value axis
value_axis = chart.value_axis
value_axis.has_major_gridlines = True
value_axis.minimum_scale = 0
value_axis.maximum_scale = 45

# Data labels
plot = chart.plots[0]
plot.has_data_labels = True
plot.data_labels.font.size = Pt(10)
plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

# Notes
slide1.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: The construction tech market is growing exponentially. "
    "From 8.5 billion in 2020 to a projected 39.4 billion by 2026."
)

print("Chart 1 (Bar) created successfully")

# ============================================================
# Chart 2: PIE CHART
# ============================================================
slide2 = prs.slides.add_slide(blank_layout)

for ph in slide2.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text_frame.text = "Technology Adoption Rate by Category"
        break

pie_data = CategoryChartData()
pie_data.categories = ['BIM', 'IoT & Sensors', 'Drones', 'AI & ML', 'Robotics', 'Cloud Platforms']
pie_data.add_series('Adoption %', (68, 45, 52, 28, 15, 72))

graphic_frame2 = slide2.shapes.add_chart(
    XL_CHART_TYPE.PIE,
    Inches(1.5), Inches(2.0), Inches(10.0), Inches(4.5),
    pie_data
)
chart2 = graphic_frame2.chart
chart2.has_legend = True
chart2.legend.position = XL_LEGEND_POSITION.RIGHT
chart2.legend.include_in_layout = False

# Data labels with percentage
plot2 = chart2.plots[0]
plot2.has_data_labels = True
plot2.data_labels.show_percentage = True
plot2.data_labels.show_category_name = False
plot2.data_labels.show_value = False
plot2.data_labels.font.size = Pt(11)
plot2.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

slide2.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: Cloud platforms and BIM lead adoption at 72% and 68% respectively."
)
print("Chart 2 (Pie) created successfully")

# ============================================================
# Chart 3: LINE CHART (with multiple series)
# ============================================================
slide3 = prs.slides.add_slide(blank_layout)

for ph in slide3.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text_frame.text = "Project Performance Improvement Over Time (%)"
        break

line_data = CategoryChartData()
line_data.categories = ['2019', '2020', '2021', '2022', '2023', '2024', '2025']
line_data.add_series('Cost Savings',  (5,  8,  12, 16, 20, 25, 30))
line_data.add_series('Schedule Reduction', (3,  6,  10, 14, 19, 24, 28))
line_data.add_series('Safety Improvement',  (2,  4,   8, 12, 18, 22, 27))

graphic_frame3 = slide3.shapes.add_chart(
    XL_CHART_TYPE.LINE_MARKERS,
    Inches(0.92), Inches(2.0), Inches(11.5), Inches(4.5),
    line_data
)
chart3 = graphic_frame3.chart
chart3.has_legend = True
chart3.legend.position = XL_LEGEND_POSITION.BOTTOM
chart3.legend.include_in_layout = False

# Data labels
plot3 = chart3.plots[0]
plot3.has_data_labels = True
plot3.data_labels.font.size = Pt(9)

# Axis styling
chart3.value_axis.has_major_gridlines = True
chart3.value_axis.minimum_scale = 0
chart3.value_axis.maximum_scale = 35

slide3.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: All three metrics show consistent improvement. "
    "Cost savings lead at 30% by 2025."
)
print("Chart 3 (Line) created successfully")

# Save
import os
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"\nSaved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
