"""Slide Master reader for Capability B (template-modifier-skill), issue #46.

Reads Slide Master-level info (master placeholders + theme) and **delegates the
full contract** to the P0 introspection engine (``template_introspector``) — no
introspection logic is duplicated.
"""
from __future__ import annotations

import logging
import sys
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation

# Reuse the P0 introspection engine from the sibling ppt-template-filler skill.
_FILLER_SCRIPTS = Path(__file__).resolve().parents[2] / "ppt-template-filler" / "scripts"
if str(_FILLER_SCRIPTS) not in sys.path:
    sys.path.insert(0, str(_FILLER_SCRIPTS))

from template_introspector import placeholder_record  # noqa: E402
from ppt_builder import get_render_contract  # noqa: E402  (US-4.1: prefer embedded JSON)

logger = logging.getLogger(__name__)


def read_master(template_path: str) -> Dict[str, Any]:
    """Read the Slide Master and return master-level + contract info.

    Returns ``{contract, master_placeholders, theme, slide_size}``. The contract
    (layouts / fingerprints / content areas) comes straight from P0
    introspection; only the master-level placeholders (not present in the
    per-layout contract) are read here.
    """
    contract = get_render_contract(template_path)
    prs = Presentation(template_path)
    master_placeholders: List[Dict[str, Any]] = []
    try:
        master = prs.slide_masters[0]
        for ph in master.placeholders:
            record = placeholder_record(ph)
            if record is not None:
                master_placeholders.append(record)
    except Exception as exc:  # pragma: no cover - defensive
        logger.warning("Slide Master read failed (%s); returning empty master", exc)

    return {
        "contract": contract,
        "master_placeholders": master_placeholders,
        "theme": contract.get("theme", {}),
        "slide_size": contract.get("slide_size", {}),
    }
