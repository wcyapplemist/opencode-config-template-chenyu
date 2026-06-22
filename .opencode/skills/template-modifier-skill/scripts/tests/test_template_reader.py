"""Tests for template_reader (Slide Master reader, issue #46)."""
from pptx import Presentation

from template_introspector import introspect
from template_reader import read_master


class TestReadMaster:
    def test_returns_contract_theme_and_slide_size(self, template_path):
        info = read_master(template_path)
        assert "contract" in info
        assert "theme" in info and info["theme"]
        assert "slide_size" in info and info["slide_size"].get("ratio")

    def test_contract_matches_fresh_introspection(self, template_path):
        """template_reader must reuse P0 introspection (no divergence)."""
        info = read_master(template_path)
        fresh = introspect(template_path)
        assert len(info["contract"]["layouts"]) == len(fresh["layouts"])
        assert info["contract"]["slide_size"] == fresh["slide_size"]

    def test_master_placeholders_is_list_of_records(self, template_path):
        info = read_master(template_path)
        mps = info["master_placeholders"]
        assert isinstance(mps, list)
        for rec in mps:
            for key in ("idx", "name", "type", "left_in", "top_in", "width_in", "height_in"):
                assert key in rec
        # Chrome placeholders must be excluded here too.
        assert all(rec["type"] not in {"FOOTER", "SLIDE_NUMBER", "DATE"} for rec in mps)
