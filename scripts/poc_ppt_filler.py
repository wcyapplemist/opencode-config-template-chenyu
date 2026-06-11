"""
poc_ppt_filler.py
=================
PPT 模板占位符探测与写入测试脚本。

功能 1：探测母版中所有幻灯片版式及其占位符信息。
功能 2：使用第一个版式创建新幻灯片，向文本占位符写入测试数据。

用法：
    pip install python-pptx
    python scripts/poc_ppt_filler.py
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

MASTER_TEMPLATE = Path(r"D:\BETEKK\master_template.pptx")
OUTPUT_FILE = Path(r"D:\BETEKK\output_test.pptx")


def analyze_placeholders(prs: Presentation) -> None:
    """遍历所有幻灯片版式，打印占位符信息。"""
    print("=" * 60)
    print("【功能 1】占位符探测 - 开始")
    print("=" * 60)

    for i, layout in enumerate(prs.slide_layouts):
        print(f"\n版式索引: {i}  |  版式名称: \"{layout.name}\"")
        print("-" * 50)

        if not layout.placeholders:
            print("  (无占位符)")
            continue

        for ph in layout.placeholders:
            ph_type = ph.placeholder_format.type
            type_name = _placeholder_type_name(ph_type)
            print(
                f"  占位符 ID: {ph.placeholder_format.idx:>3}  |  "
                f"类型: {type_name:<25s}  |  "
                f"名称: \"{ph.name}\""
            )

    print("\n" + "=" * 60)
    print("【功能 1】占位符探测 - 完成")
    print("=" * 60)


def write_test_slide(prs: Presentation) -> None:
    """使用第一个版式创建幻灯片，向文本占位符写入测试数据。"""
    print("\n" + "=" * 60)
    print("【功能 2】写入测试 - 开始")
    print("=" * 60)

    layout = prs.slide_layouts[0]
    print(f"使用版式: \"{layout.name}\" (索引 0)")

    slide = prs.slides.add_slide(layout)

    test_messages = [
        "BETEKK POC Test",
        "Hello World",
        "这是第三个占位符的测试内容",
        "第四个备用文本",
        "第五个备用文本",
    ]
    msg_index = 0

    for ph in slide.placeholders:
        ph_type = ph.placeholder_format.type

        if ph_type in (
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
            PP_PLACEHOLDER.SUBTITLE,
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
        ):
            text = test_messages[msg_index % len(test_messages)]
            ph.text = text
            print(f"  写入占位符 ID {ph.placeholder_format.idx} ({ph.name}): \"{text}\"")
            msg_index += 1
        else:
            type_name = _placeholder_type_name(ph_type)
            print(f"  跳过占位符 ID {ph.placeholder_format.idx} ({type_name})")

    if msg_index == 0:
        print("  警告: 未找到任何文本类占位符，写入跳过。")

    prs.save(str(OUTPUT_FILE))
    print(f"\n测试文件已保存至: {OUTPUT_FILE}")
    print("=" * 60)
    print("【功能 2】写入测试 - 完成")
    print("=" * 60)


def _placeholder_type_name(ph_type) -> str:
    """将占位符枚举值转为可读字符串。"""
    try:
        return str(ph_type).split(".")[-1]
    except Exception:
        return str(ph_type)


def main() -> None:
    if not MASTER_TEMPLATE.exists():
        print(f"错误: 找不到母版文件 {MASTER_TEMPLATE}")
        sys.exit(1)

    print(f"加载母版文件: {MASTER_TEMPLATE}")
    prs = Presentation(str(MASTER_TEMPLATE))
    print(f"母版包含 {len(prs.slide_layouts)} 个版式\n")

    analyze_placeholders(prs)
    write_test_slide(prs)

    print("\n全部流程执行完毕。")


if __name__ == "__main__":
    main()
